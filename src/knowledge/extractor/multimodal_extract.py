import sys
import json
import os
import urllib.request
import urllib.error
import base64
import time
from io import BytesIO

# Ensure opt/homebrew/bin is in PATH for subprocesses (like pdftoppm / tesseract)
os.environ['PATH'] = '/opt/homebrew/bin:/usr/local/bin:/usr/bin:/bin:/usr/sbin:/sbin:' + os.environ.get('PATH', '')

try:
    import pytesseract
    pytesseract.pytesseract.tesseract_cmd = '/opt/homebrew/bin/tesseract'
except ImportError:
    pass

# Strict JSON Schema for Gemini REST Structured Output
RESPONSE_SCHEMA = {
    "type": "OBJECT",
    "properties": {
        "figures": {
            "type": "ARRAY",
            "description": "List of visual graphics/figures found on the page. Empty if none.",
            "items": {
                "type": "OBJECT",
                "required": ["id", "archetype", "caption", "details"],
                "properties": {
                    "id": { "type": "STRING", "description": "Unique ID for the figure, e.g. FIG_09_01" },
                    "archetype": {
                        "type": "STRING",
                        "enum": [
                            "TIME-COURSE & PATHOPHYSIOLOGICAL TRENDS",
                            "SUBCELLULAR BIOCHEMICAL CASCADE DIAGRAMS",
                            "COORDINATE X-Y GRAPHS & COMPLEMENTARY PANELS",
                            "ANATOMICAL NETWORKS & MICROCIRCUIT MAPS",
                            "PHYSIOLOGICAL WAVEFORMS & TRACINGS"
                        ]
                    },
                    "caption": { "type": "STRING", "description": "High-fidelity caption or text description of the figure" },
                    "details": {
                        "type": "OBJECT",
                        "description": "Specific details for the archetype. Only include the fields corresponding to the archetype class.",
                        "properties": {
                            "time_domain_x_axis": {
                                "type": "OBJECT",
                                "properties": {
                                    "unit": { "type": "STRING" },
                                    "range": { "type": "STRING" }
                                }
                            },
                            "y_axis_polarity": {
                                "type": "OBJECT",
                                "properties": {
                                    "upper_hemisphere_label": { "type": "STRING" },
                                    "lower_hemisphere_label": { "type": "STRING" }
                                }
                            },
                            "curves": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "label_or_color": { "type": "STRING" },
                                        "initiate_epoch": { "type": "STRING" },
                                        "peak_or_trough_epoch": { "type": "STRING" },
                                        "decay_duration": { "type": "STRING" }
                                    }
                                }
                            },
                            "molecular_markers": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "marker_name": { "type": "STRING" },
                                        "curve_segment": { "type": "STRING" },
                                        "time_window": { "type": "STRING" }
                                    }
                                }
                            },
                            "temporal_alignments": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "epoch": { "type": "STRING" },
                                        "event_description": { "type": "STRING" },
                                        "markers_active": { "type": "ARRAY", "items": { "type": "STRING" } }
                                    }
                                }
                            },
                            "compartments": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "name": { "type": "STRING" },
                                        "type": { "type": "STRING" },
                                        "boundaries": { "type": "ARRAY", "items": { "type": "STRING" } }
                                    }
                                }
                            },
                            "transmembrane_elements": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "name": { "type": "STRING" },
                                        "compartment_boundary": { "type": "STRING" },
                                        "activity_type": { "type": "STRING" }
                                    }
                                }
                            },
                            "internal_chemical_transactions": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "sequence_string": { "type": "STRING" },
                                        "substrate": { "type": "STRING" },
                                        "enzyme_or_catalyst": { "type": "STRING" },
                                        "product": { "type": "STRING" },
                                        "compartment": { "type": "STRING" }
                                    }
                                }
                            },
                            "cross_compartment_pathways": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "source_compartment": { "type": "STRING" },
                                        "target_compartment": { "type": "STRING" },
                                        "signaling_molecule": { "type": "STRING" },
                                        "action_or_receptor": { "type": "STRING" }
                                    }
                                }
                            },
                            "downstream_outputs": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "functional_output": { "type": "STRING" },
                                        "trigger_cascade": { "type": "STRING" }
                                    }
                                }
                            },
                            "panels": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "panel_id": { "type": "STRING" },
                                        "axes": {
                                            "type": "OBJECT",
                                            "properties": {
                                                "x_axis": { "type": "OBJECT", "properties": { "label": { "type": "STRING" }, "unit": { "type": "STRING" } } },
                                                "y_axis": { "type": "OBJECT", "properties": { "label": { "type": "STRING" }, "unit": { "type": "STRING" } } }
                                            }
                                        },
                                        "mathematical_nature": { "type": "STRING" },
                                        "curves": {
                                            "type": "ARRAY",
                                            "items": {
                                                "type": "OBJECT",
                                                "properties": {
                                                    "legend_label": { "type": "STRING" },
                                                    "characteristic_nature": { "type": "STRING" }
                                                }
                                            }
                                        },
                                        "coordinate_inflections": {
                                            "type": "ARRAY",
                                            "items": {
                                                "type": "OBJECT",
                                                "properties": {
                                                    "x_threshold": { "type": "STRING" },
                                                    "y_value_or_change": { "type": "STRING" },
                                                    "inflection_description": { "type": "STRING" }
                                                }
                                            }
                                        }
                                    }
                                },
                            },
                            "nodes_list": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "id": { "type": "STRING" },
                                        "name": { "type": "STRING" },
                                        "relational_variants": { "type": "ARRAY", "items": { "type": "STRING" } },
                                        "prevalence_percentage": { "type": "NUMBER" }
                                    }
                                }
                            },
                            "source_target_vectors": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "source": { "type": "STRING" },
                                        "target": { "type": "STRING" },
                                        "direction": { "type": "STRING", "enum": ["Unidirectional", "Bidirectional", "Feed-forward loop"] },
                                        "sign": { "type": "STRING", "enum": ["Excitatory", "Inhibitory"] },
                                        "neurotransmitter": { "type": "STRING" }
                                    }
                                }
                            },
                            "channels": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "source_marker": { "type": "STRING" },
                                        "amplitude_range": { "type": "STRING" }
                                    }
                                }
                            },
                            "transient_synchronous_events": {
                                "type": "ARRAY",
                                "items": {
                                    "type": "OBJECT",
                                    "properties": {
                                        "time_window": { "type": "STRING" },
                                        "flags_or_annotations": { "type": "ARRAY", "items": { "type": "STRING" } },
                                        "morphology_pattern": { "type": "STRING" },
                                        "channel_relationships": { "type": "STRING" }
                                    }
                                }
                            },
                            "ecg_findings": {
                                "type": "OBJECT",
                                "description": "Only populate for ECG/EKG tracings. Omit entirely for non-ECG waveforms.",
                                "properties": {
                                    "lead_shown": { "type": "STRING", "description": "e.g. Lead II, V1, limb leads I/II/III, or 'not specified'" },
                                    "named_rhythm": { "type": "STRING", "description": "e.g. complete heart block, asystole, ventricular fibrillation, normal sinus rhythm, atrial fibrillation" },
                                    "rate_bpm": { "type": "STRING", "description": "Heart rate only if depicted/labeled or directly calculable from labeled calibration marks" },
                                    "pr_interval_ms": { "type": "STRING" },
                                    "qrs_duration_ms": { "type": "STRING" },
                                    "qt_interval_ms": { "type": "STRING" },
                                    "st_segment_finding": { "type": "STRING", "description": "e.g. ST elevation lead II, ST depression scooped/digoxin-type, isoelectric" },
                                    "t_wave_finding": { "type": "STRING", "description": "e.g. tented T waves, T wave inversion, normal" },
                                    "other_morphology": { "type": "STRING", "description": "e.g. RBBB rabbit-ears pattern, U waves present, f-waves (AFib)" }
                                }
                            },
                            "capnography_findings": {
                                "type": "OBJECT",
                                "description": "Only populate for capnography/EtCO2 tracings. Omit entirely for non-capnography waveforms.",
                                "properties": {
                                    "named_pattern": { "type": "STRING", "description": "e.g. obstructive shark-fin upstroke, curare cleft, cardiogenic oscillations, rebreathing elevated baseline, esophageal intubation flatline, cuff leak dilution, bucking artifact, normal" },
                                    "baseline_etco2_value": { "type": "STRING" },
                                    "plateau_etco2_value": { "type": "STRING" },
                                    "waveform_phase_described": { "type": "STRING", "description": "e.g. phase II upstroke, phase III plateau, inspiratory downstroke" }
                                }
                            }
                        }
                    }
                }
            }
        }
    }
}

def load_api_key():
    # Try environment variable first
    api_key = os.environ.get('GEMINI_API_KEY')
    if api_key:
        return api_key
    # Fallback: search for .env in current and parent directories
    curr_dir = os.path.dirname(os.path.abspath(__file__))
    for _ in range(4):
        env_path = os.path.join(curr_dir, '.env')
        if os.path.exists(env_path):
            with open(env_path, 'r') as f:
                for line in f:
                    if line.strip().startswith('GEMINI_API_KEY='):
                        return line.split('=', 1)[1].strip()
        curr_dir = os.path.dirname(curr_dir)
    return None

def call_gemini_api(img_base64, api_key, model_name="gemini-2.5-flash"):
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
    
    prompt = """List only the visible text and structural elements explicitly present in this image. Do not extrapolate, assume, or provide external medical knowledge. If a connection is not explicitly drawn, do not include it.

You are a high-fidelity visual extraction engine for an advanced medical document parser. 
Your task is to analyze the provided page image and extract all visual elements (graphics, charts, schematics, and waveforms) into a standardized, fully semantic JSON structure optimized for computational simulation models.

For every visual element identified on the page, you must class-map the graphic to one of the following five archetypes and parse it according to the specific rules specified below.

---

### ARCHETYPE 1: TIME-COURSE & PATHOPHYSIOLOGICAL TRENDS
- **Detection Criteria:** Line charts or timeline graphs tracking injury, progression, or drug effects across units of time.
- **Extraction Protocol:**
  1. Extract the explicit units of the Time Domain X-axis (e.g., Minutes, Hours, Days, Weeks).
  2. Identify the Y-axis polarity (e.g., Upper hemisphere = Injury/Destruction, Lower hemisphere = Protection/Repair).
  3. Identify each discrete trace/curve by its text-label or color profile.
  4. Trace the timeline behavior of each curve: determine the precise time-epoch where it initiates, where it peaks/troughs, and its decay duration.
  5. Map specific molecular markers (e.g., Caspases, Glutamate, IL-10) directly to the specific curve segment and time window where they are clustered.
  6. Fill in 'temporal_alignments' to represent the epochs, event descriptions, and active markers.

### ARCHETYPE 2: SUBCELLULAR BIOCHEMICAL CASCADE DIAGRAMS
- **Detection Criteria:** Cell-surface diagrams containing receptors, second-messenger pathways, and cross-compartmental signaling vectors.
- **Extraction Protocol:**
  1. Identify and define the explicit boundaries and Compartments (e.g., "Presynaptic Neuron Terminal", "Astrocyte Foot", "Vascular Endothelium").
  2. Map all transmembrane receptors and channels (e.g., NMDAR, mGluR) to their specific compartment boundaries.
  3. Convert all internal chemical transactions into sequence strings: [Substrate_A] -> [Enzyme/Catalyst] -> [Product_B].
  4. Map cross-compartment pathways: Identify the signaling molecule released from Compartment X and its destination receptor/action on Compartment Y.
  5. Capture downstream functional outputs (e.g., "Dilation", "Constriction").

### ARCHETYPE 3: COORDINATE X-Y GRAPHS & COMPLEMENTARY PANELS
- **Detection Criteria:** Multi-panel plots (labeled A, B, C, etc.) describing physiological metrics across absolute values, percentages, or saturation rates.
- **Extraction Protocol:**
  1. Create a distinct data object for each panel character (A, B, C).
  2. Identify the exact string labels and measurement units for both the X and Y axes.
  3. Determine the structural mathematical nature of each line curve: "Inversely Linear", "Sigmoidal/Logarithmic Plateau", "Exponential Decay", "Linear", or "Other".
  4. Extract critical turning points (coordinate_inflections): identify the X-axis inflection threshold where a curve rapidly changes its slope (e.g., "CBF remains stable until PaO2 drops below 60 mmHg, where an exponential increase occurs").
  5. Map legend labels explicitly to their corresponding data lines.

### ARCHETYPE 4: ANATOMICAL NETWORKS & MICROCIRCUIT MAPS
- **Detection Criteria:** Box-and-arrow diagrams, neuroanatomical circuit schematics, or vascular trees (e.g., Circle of Willis).
- **Extraction Protocol:**
  1. Catalog all anatomical node identities into an entity dictionary (as a list of nodes under 'nodes_list' containing: 'id', 'name', 'relational_variants', and 'prevalence_percentage').
  2. Map relational variants or statistical occurrences if present (e.g., Anatomical Variant X = 27% prevalence).
  3. Trace every line vector between nodes and extract edge attributes (source_target_vectors):
     - "direction": Unidirectional, Bidirectional, or Feed-forward loop.
     - "sign": Excitatory (+ / Glutamatergic / arrowhead) vs. Inhibitory (- / GABAergic / T-bar perpendicular terminal).
     - Neurotransmitter and sign details if present.

### ARCHETYPE 5: PHYSIOLOGICAL WAVEFORMS & TRACINGS
- **Detection Criteria:** Parallel multichannel continuous waveforms (e.g., polysomnography, EEG, arterial line traces), or a single-channel clinical monitor tracing (ECG strip, capnography/EtCO2 trace).
- **Extraction Protocol:**
  1. Segment the visualization into parallel vertical channels and extract the exact source marker of each channel (e.g., SaO2, Resp Flow, Thorax, Abdomen).
  2. Detect transient synchronous events: match localized bounding boxes or annotation flags (e.g., "CH", "OA") to the exact time-window slice on the trace timeline.
  3. Describe the wave morphological pattern during events (e.g., "Crescendo-decrescendo modulation", "Flatline/Cessation of flow", "Phase paradox between Thorax and Abdomen channels").
  4. If the tracing is specifically an ECG/EKG strip, additionally populate `ecg_findings` with the lead(s) shown, the named rhythm, the rate if depicted or directly calculable from labeled calibration marks, any labeled interval values (PR/QRS/QT), and ST-segment/T-wave findings. Read any digital numeric readout shown alongside the trace (e.g. a monitor's displayed HR) as a separate fact from what the trace itself actually shows — do not assume they agree; a mismatch between a monitor's numeric display and the true waveform is sometimes the entire clinical point of the figure.
  5. If the tracing is specifically a capnography/EtCO2 trace, additionally populate `capnography_findings` with the named morphological pattern (shark-fin, curare cleft, cardiogenic oscillations, rebreathing, esophageal, cuff-leak, bucking, or normal) if one is clinically depicted, plus baseline/plateau values if labeled.
  6. Leave `ecg_findings` and `capnography_findings` entirely absent for EEG/polysomnography/arterial-line/other non-cardiorespiratory waveforms.

---

### CRITICAL SPATIAL PROXIMITY RULE FOR NETWORKS (ARCHETYPE 4):
When tracing structural edges in box-and-arrow diagrams or circuit schematics, you MUST NOT map a source node to a target node simply because they exist on the same diagram page. You MUST trace each vector line continuously with pixel-level tracking from its explicit origin point to its explicit point of termination (arrowhead or T-bar perpendicular terminal). Specifically:
- If an axon or projection line terminates locally inside a coordinate structure (e.g., local cortical inhibition between two adjacent nodes), do NOT map it to remote subcortical structures located further down the coordinate space.
- Only register a source_target_vector if you can visually confirm a continuous drawn line or arrow connecting the two nodes.
- Nodes that are merely co-located on the same page without a visible connecting line are NOT edges.
- When in doubt about whether a line connects two specific nodes, do NOT include the edge. Precision over recall.

### IMPORTANT RULES:
1. Do not let downstream numeric strings drift into an unassigned text state; every text token found inside a graphic region must be assigned to an attribute key or a relational node parent.
2. If the page does not contain any visual graphics/figures, return an empty array for "figures".
3. Return a high-fidelity semantic parsing result.
"""

    request_body = {
        "contents": [
            {
                "parts": [
                    {
                        "text": prompt
                    },
                    {
                        "inlineData": {
                            "mimeType": "image/jpeg",
                            "data": img_base64
                        }
                    }
                ]
            }
        ],
        "generationConfig": {
            "responseMimeType": "application/json",
            "responseSchema": RESPONSE_SCHEMA,
            "temperature": 0.1
        }
    }
    
    request_data = json.dumps(request_body).encode('utf-8')
    headers = {'Content-Type': 'application/json'}
    req = urllib.request.Request(url, data=request_data, headers=headers, method='POST')
    
    max_retries = 4
    delay = 3
    for attempt in range(1, max_retries + 1):
        try:
            with urllib.request.urlopen(req, timeout=120) as response:
                res_data = json.loads(response.read().decode('utf-8'))
                candidate = res_data.get('candidates', [{}])[0]
                text = candidate.get('content', {}).get('parts', [{}])[0].get('text', '{}')
                parsed = json.loads(text)
                
                # Transform nodes_list to nodes dictionary to match types perfectly!
                figures = parsed.get('figures', [])
                for fig in figures:
                    details = fig.get('details', {})
                    if 'nodes_list' in details:
                        nodes_list = details.get('nodes_list') or []
                        nodes_dict = {}
                        for n in nodes_list:
                            n_id = n.get('id') or n.get('name')
                            if n_id:
                                nodes_dict[n_id] = {
                                    "name": n.get('name') or n_id,
                                    "relational_variants": n.get('relational_variants') or [],
                                    "prevalence_percentage": n.get('prevalence_percentage')
                                }
                        details['nodes'] = nodes_dict
                        del details['nodes_list']
                        
                return figures
        except Exception as e:
            if attempt == max_retries:
                print(f"  [GEMINI ERROR] API call failed after {max_retries} attempts: {str(e)}", file=sys.stderr)
                return []
            print(f"  [GEMINI WARNING] Attempt {attempt} failed, retrying in {delay}s: {str(e)}", file=sys.stderr)
            time.sleep(delay)
            delay *= 2
            
    return []

def decode_unicode_escapes(text):
    import re
    if not text:
        return ""
    # Matches both /uXXXX and \uXXXX escape patterns
    pattern = re.compile(r'[/\\]u([0-9a-fA-F]{4})')
    def replace(match):
        try:
            return chr(int(match.group(1), 16))
        except ValueError:
            return match.group(0)
    decoded = pattern.sub(replace, text)
    
    # Matches any non-printable ASCII control character (excluding \t, \n, \r)
    control_char_pat = re.compile(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]')
    lines = decoded.split('\n')
    cleaned_lines = [line for line in lines if not control_char_pat.search(line)]
    return '\n'.join(cleaned_lines)

def extract_pdf_text(file_path):
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        text = decode_unicode_escapes(text)
        pages.append({
            "page_number": i + 1,
            "text": text,
            "char_count": len(text)
        })
    return {
        "total_pages": len(reader.pages),
        "pages": pages
    }

def extract_pdf_ocr(file_path):
    from pdf2image import convert_from_path
    import pytesseract
    
    # Use eng_best (tessdata_best LSTM model) for maximum accuracy on medical text
    # DPI 300 to capture subscripts, superscripts, and small medical notation
    lang = 'eng_best' if os.path.exists('/opt/homebrew/share/tessdata/eng_best.traineddata') else 'eng'
    print(f"  [OCR] Using tessdata model: {lang} at 300 DPI", file=sys.stderr)
    
    # Convert PDF to images at higher DPI for medical text accuracy
    images = convert_from_path(file_path, dpi=300)
    pages = []
    for i, img in enumerate(images):
        # PSM 6: Assume a single uniform block of text (best for full pages)
        custom_config = f'--oem 1 --psm 6 -l {lang}'
        text = pytesseract.image_to_string(img, config=custom_config) or ""
        text = decode_unicode_escapes(text)
        pages.append({
            "page_number": i + 1,
            "text": text,
            "char_count": len(text)
        })
    return {
        "total_pages": len(images),
        "pages": pages
    }

def extract_pdf_visual(file_path):
    from pypdf import PdfReader
    from pdf2image import convert_from_path
    
    api_key = load_api_key()
    if not api_key:
        raise ValueError("GEMINI_API_KEY not found in environment or .env file.")
        
    model_name = os.environ.get('GEMINI_MODEL', 'gemini-2.5-flash')
    print(f"  [VISUAL PARSE] Reading PDF to analyze page contents...", file=sys.stderr)
    reader = PdfReader(file_path)
    
    pages = []
    
    # Pre-filter figure keywords to save massive rendering time and API costs
    fig_keywords = ["fig.", "figure", "diagram", "schematic", "waveform", "chart", "graph", "timeline"]
    
    for i, page in enumerate(reader.pages):
        page_num = i + 1
        text = (page.extract_text() or "").lower()
        
        has_fig = any(kw in text for kw in fig_keywords)
        
        if not has_fig:
            # Skip page, it does not contain figures
            pages.append({
                "page_number": page_num,
                "figures": []
            })
            continue
            
        print(f"  [VISUAL PARSE] Figure keywords detected on page {page_num}. Converting to image...", file=sys.stderr)
        try:
            # Convert ONLY the specific page to image to keep performance lightning fast!
            images = convert_from_path(file_path, first_page=page_num, last_page=page_num, dpi=150)
            if not images:
                pages.append({
                    "page_number": page_num,
                    "figures": []
                })
                continue
                
            img = images[0]
            print(f"  [VISUAL PARSE] Processing page {page_num} via Gemini ({model_name})...", file=sys.stderr)
            
            # Convert image to base64 JPEG bytes
            buffered = BytesIO()
            img.save(buffered, format="JPEG")
            img_str = base64.b64encode(buffered.getvalue()).decode('utf-8')
            
            # Call Gemini API
            figures = call_gemini_api(img_str, api_key, model_name)
            
            pages.append({
                "page_number": page_num,
                "figures": figures
            })
        except Exception as ex:
            print(f"  [VISUAL PARSE ERROR] Page {page_num} failed to parse: {str(ex)}", file=sys.stderr)
            pages.append({
                "page_number": page_num,
                "figures": []
            })
        
    return {
        "total_pages": len(reader.pages),
        "pages": pages
    }

def extract_image_ocr(file_path):
    from PIL import Image
    import pytesseract
    
    # Use eng_best (tessdata_best LSTM model) for maximum accuracy on medical text
    lang = 'eng_best' if os.path.exists('/opt/homebrew/share/tessdata/eng_best.traineddata') else 'eng'
    print(f"  [OCR] Using tessdata model: {lang}", file=sys.stderr)
    
    img = Image.open(file_path)
    # PSM 6: Assume a single uniform block of text
    custom_config = f'--oem 1 --psm 6 -l {lang}'
    text = pytesseract.image_to_string(img, config=custom_config) or ""
    text = decode_unicode_escapes(text)
    return {
        "total_pages": 1,
        "pages": [{
            "page_number": 1,
            "text": text,
            "char_count": len(text)
        }]
    }

def extract_image_visual(file_path):
    from PIL import Image
    
    api_key = load_api_key()
    if not api_key:
        raise ValueError("GEMINI_API_KEY not found in environment or .env file.")
        
    model_name = os.environ.get('GEMINI_MODEL', 'gemini-2.5-flash')
    
    img = Image.open(file_path)
    # Convert image to base64 JPEG bytes
    buffered = BytesIO()
    # Convert to RGB if needed to save as JPEG
    if img.mode in ('RGBA', 'LA'):
        img = img.convert('RGB')
    img.save(buffered, format="JPEG")
    img_str = base64.b64encode(buffered.getvalue()).decode('utf-8')
    
    print(f"  [VISUAL PARSE] Processing standalone image via Gemini ({model_name})...", file=sys.stderr)
    figures = call_gemini_api(img_str, api_key, model_name)
    
    return {
        "total_pages": 1,
        "pages": [{
            "page_number": 1,
            "figures": figures
        }]
    }

def extract_pptx_text(file_path):
    from pptx import Presentation
    
    prs = Presentation(file_path)
    pages = []
    
    for i, slide in enumerate(prs.slides):
        slide_parts = []
        # Get slide title if exists
        if slide.shapes.title:
            slide_parts.append(slide.shapes.title.text)
            
        # Get all text from body shapes
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        slide_parts.append(paragraph.text)
                        
        slide_text = "\n".join(slide_parts)
        slide_text = decode_unicode_escapes(slide_text)
        pages.append({
            "page_number": i + 1,
            "text": slide_text,
            "char_count": len(slide_text)
        })
        
    return {
        "total_pages": len(prs.slides),
        "pages": pages
    }

def main():
    if len(sys.argv) < 3:
        print(json.dumps({"error": "Usage: python3 multimodal_extract.py [mode] [file_path]"}))
        sys.exit(1)
        
    mode = sys.argv[1]
    file_path = sys.argv[2]
    
    if not os.path.exists(file_path):
        print(json.dumps({"error": f"File not found: {file_path}"}))
        sys.exit(1)
        
    try:
        if mode == "pdf_text":
            result = extract_pdf_text(file_path)
        elif mode == "pdf_ocr":
            result = extract_pdf_ocr(file_path)
        elif mode == "pdf_visual":
            result = extract_pdf_visual(file_path)
        elif mode == "image":
            result = extract_image_ocr(file_path)
        elif mode == "image_visual":
            result = extract_image_visual(file_path)
        elif mode == "pptx":
            result = extract_pptx_text(file_path)
        else:
            result = {"error": f"Unknown mode: {mode}"}
            
        print(json.dumps(result))
    except Exception as e:
        print(json.dumps({"error": str(e)}))
        sys.exit(1)

if __name__ == "__main__":
    main()
