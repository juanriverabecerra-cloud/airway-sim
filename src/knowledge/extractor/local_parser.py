import sys
import json
import os
import re
import math
import shutil
import traceback
import unicodedata

try:
    import fitz  # PyMuPDF
except ImportError:
    print(json.dumps({"error": "PyMuPDF (fitz) is not installed."}))
    sys.exit(1)

try:
    import pytesseract
    # Resolve tesseract from PATH first so this runs on Linux/CI/other machines;
    # only fall back to the Homebrew-on-macOS path this was originally pinned to.
    _tesseract_path = shutil.which('tesseract') or '/opt/homebrew/bin/tesseract'
    if os.path.exists(_tesseract_path):
        pytesseract.pytesseract.tesseract_cmd = _tesseract_path
except ImportError:
    pass

try:
    from PIL import Image
except ImportError:
    pass

# silences fitz C-level warnings to ensure stdout remains pure JSON
try:
    fitz.tools.set_warning_level(0)
except AttributeError:
    pass

# Scratch directory for cropped figures/tables, anchored to the repo root (not CWD)
# so output location is stable regardless of which directory spawns this script.
_REPO_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..', '..'))
SCRATCH_DIR = os.environ.get('LOCAL_PARSER_SCRATCH_DIR') or os.path.join(_REPO_ROOT, 'scratch', 'extracted_images')

# Neurotransmitter/keyword cues used to deterministically infer connection polarity
# from text that is actually present in the figure (never guessed/fabricated).
_INHIBITORY_CUES = ['gaba', 'glycine', 'inhibitory', '(-)', '(–)']
_EXCITATORY_CUES = ['glutamate', 'glutamatergic', 'excitatory', '(+)']

def sanitize_fig_id(cap_text, xref, page_num):
    # Try to find standard figure identifier like "Fig. 9.1"
    fig_match = re.search(r'(?:Fig\.|Figure)\s*([\d\.\-_]+)', cap_text, re.IGNORECASE)
    if fig_match:
        clean_num = fig_match.group(1).replace('.', '_').replace('-', '_').strip('_')
        return f"FIG_{clean_num}"
    return f"FIG_PAGE_{page_num}_xref_{xref}"

def is_real_heading(first_line):
    """
    Distinguishes an actual section heading from ordinary paragraph text or a
    stray isolated label (e.g. a figure's axis/legend word floating in its own
    text block). The previous heuristic ("starts with any capitalized word")
    matched almost every sentence and every short capitalized label, which
    fragmented each page into ~12 mostly-empty "sections" instead of a few
    real ones — same total text, but 10x as many tiny database rows, which is
    what was bloating both the database and the search index. Tightening this
    doesn't drop any text: a line that fails this check still gets kept, just
    appended to the current section's body instead of starting a new one.
    """
    first_line = first_line.strip()
    if not first_line or len(first_line) > 100:
        return False

    # Strongest signal: a true all-caps heading line (Miller's convention).
    if first_line.isupper() and len(first_line) >= 4:
        return True

    # Numbered section header, e.g. "19.1 Effects of Anesthetics" or "1. Overview"
    if re.match(r'^\d+(?:\.\d+)*\.?\s+[A-Z]', first_line):
        return True

    # Genuine multi-word Title Case heading: most words capitalized, and not
    # ending in sentence punctuation (a real sentence almost always does).
    if first_line[-1] not in '.,:;':
        words = first_line.split()
        if len(words) >= 2:
            cap_ratio = sum(1 for w in words if w[:1].isupper()) / len(words)
            if cap_ratio >= 0.7:
                return True

    return False

def get_closest_heading(blocks, cap_block_idx):
    # Find the nearest preceding heading block
    for i in range(cap_block_idx - 1, -1, -1):
        block = blocks[i]
        if block[6] == 0:  # Text block
            text = block[4].strip()
            lines = [l.strip() for l in text.split('\n') if l.strip()]
            if lines and is_real_heading(lines[0]):
                return lines[0]
    return ""

def closest_heading_before_y(blocks, y):
    """Finds the nearest heading preceding a given y-coordinate (used for tables,
    which aren't anchored to a caption block index the way figures are)."""
    idx = 0
    for i, b in enumerate(blocks):
        if b[6] == 0 and b[1] <= y:
            idx = i + 1
    return get_closest_heading(blocks, idx)

def rect_overlap_ratio(rect_a, rect_b):
    """Intersection area as a fraction of the smaller rect's area. Used to detect
    when a flowing-text block actually belongs to a table that was already
    captured structurally, so it can be excluded from the prose pass."""
    inter = rect_a & rect_b
    if inter.is_empty:
        return 0.0
    inter_area = inter.width * inter.height
    area_a = rect_a.width * rect_a.height
    area_b = rect_b.width * rect_b.height
    smaller = min(area_a, area_b)
    return inter_area / smaller if smaller > 0 else 0.0

def find_nearby_caption(blocks, target_rect, prefix_regex, max_dist=60):
    """Finds the closest text block whose first line matches prefix_regex
    (e.g. 'Table 19-2') within max_dist points of the target rect's edges."""
    best_text = ""
    best_dist = float('inf')
    for block in blocks:
        if block[6] != 0:
            continue
        text = block[4].strip()
        if not text:
            continue
        first_line = text.split('\n')[0].strip()
        if re.match(prefix_regex, first_line, re.IGNORECASE):
            bbox = block[:4]
            dist = min(abs(bbox[3] - target_rect.y0), abs(bbox[1] - target_rect.y1))
            if dist < best_dist and dist <= max_dist:
                best_dist = dist
                best_text = first_line
    return best_text

_NEUROTRANSMITTER_CUES = {
    'gabaergic': 'GABA', 'gaba': 'GABA',
    'glutamatergic': 'Glutamate', 'glutamate': 'Glutamate',
    'glycinergic': 'Glycine', 'glycine': 'Glycine',
    'dopaminergic': 'Dopamine', 'dopamine': 'Dopamine',
    'serotonergic': 'Serotonin', 'serotonin': 'Serotonin', '5-ht': 'Serotonin',
    'noradrenergic': 'Norepinephrine', 'norepinephrine': 'Norepinephrine',
    'cholinergic': 'Acetylcholine', 'acetylcholine': 'Acetylcholine',
    'histaminergic': 'Histamine', 'histamine': 'Histamine',
    'orexinergic': 'Orexin', 'orexin': 'Orexin',
}

def infer_sign_and_neurotransmitter(text_a, text_b, extra_text=''):
    """
    Infers connection sign and neurotransmitter identity ONLY from explicit
    textual cues physically present in the figure (node labels, or text
    annotated along the connector path itself via `extra_text`) — never from
    drawing style (line color/dash pattern are illustrator/publisher-specific
    conventions, not a reliable universal signal across arbitrary source
    textbooks). Consistent with the zero-hallucination design principle: a
    detected connector with no such cue is reported with 'sign'/
    'neurotransmitter' omitted rather than defaulted to a guess.
    """
    combined = f"{text_a} {text_b} {extra_text}".lower()
    result = {}
    if any(cue in combined for cue in _INHIBITORY_CUES):
        result["sign"] = "Inhibitory"
    elif any(cue in combined for cue in _EXCITATORY_CUES):
        result["sign"] = "Excitatory"
    for cue, name in _NEUROTRANSMITTER_CUES.items():
        if cue in combined:
            result["neurotransmitter"] = name
            break
    return result

def group_words_into_phrases(words):
    """
    Shared phrase-clustering logic: merges adjacent word boxes on the same line
    into phrases. Works identically whether the words came from OCR (uncertain,
    needs a confidence score) or directly from the PDF's own text layer (exact,
    always full-confidence). Keeping one implementation means native-text and
    OCR'd figures get identical downstream handling.
    """
    words = sorted(words, key=lambda w: (w['y0'], w['x0']))
    grouped = []
    for w in words:
        merged = False
        for g in grouped:
            # Group adjacent words on same line (y diff < 15, horizontal gap < 30)
            if abs(w['y0'] - g['y0']) < 15 and abs(w['x0'] - g['x1']) < 30:
                g['text'] += ' ' + w['text']
                g['x1'] = max(g['x1'], w['x1'])
                g['y0'] = min(g['y0'], w['y0'])
                g['y1'] = max(g['y1'], w['y1'])
                g['conf'] = (g['conf'] + w['conf']) / 2
                merged = True
                break
        if not merged:
            grouped.append(dict(w))

    result_boxes = []
    for g in grouped:
        text = g['text'].strip()
        if g['conf'] < 75:
            text += " [uncertain]"

        # Keep clean verbatim characters without guessing
        if len(text.replace(" [uncertain]", "").strip()) > 1:
            result_boxes.append({
                "text": text,
                "bbox": [g['x0'], g['y0'], g['x1'], g['y1']],
                "confidence": round(g['conf'], 1)
            })
    return result_boxes

def cluster_ocr_phrases(img_path):
    lang = 'eng_best' if os.path.exists('/opt/homebrew/share/tessdata/eng_best.traineddata') else 'eng'
    img = Image.open(img_path)
    # PSM 11: Sparse text with coordinates
    data = pytesseract.image_to_data(img, config=f'--oem 1 --psm 11 -l {lang}', output_type=pytesseract.Output.DICT)

    words = []
    for i in range(len(data['text'])):
        w = data['text'][i].strip()
        try:
            c = float(data['conf'][i])
        except (ValueError, TypeError):
            c = 0.0
        if w and c > 0:
            words.append({
                'text': w,
                'x0': data['left'][i],
                'y0': data['top'][i],
                'x1': data['left'][i] + data['width'][i],
                'y1': data['top'][i] + data['height'][i],
                'conf': c
            })

    return group_words_into_phrases(words)

_RUNNING_HEADER_RE = re.compile(r"^SECTION\s+[IVXLC0-9]+\s*[•·].+", re.IGNORECASE)

def strip_furniture_phrases(phrases, furniture=None):
    """Drop figure-label phrases that are actually page furniture — a bare page
    number, or a running header caught inside a figure's crop rect when the figure
    sits against the top/bottom margin (e.g. '356', 'SECTION II • Anesthetic
    Physiology' leaking into a graph's labels). Real figure labels are untouched."""
    out = []
    for p in phrases or []:
        t = (p.get("text") or "").replace(" [uncertain]", "").strip()
        if re.fullmatch(r"\d{1,4}", t):
            continue
        if _RUNNING_HEADER_RE.match(t):
            continue
        if furniture and _furniture_sig(t) in furniture:
            continue
        out.append(p)
    return out

def _graph_val(t):
    return float(t) if re.fullmatch(r'-?\d+(?:\.\d+)?', t) else None

def _graph_monotonic(vals):
    if len(vals) < 3:
        return False
    return all(b > a for a, b in zip(vals, vals[1:])) or all(b < a for a, b in zip(vals, vals[1:]))

def _graph_scale(vals):
    v = sorted(vals)
    if len(v) < 3:
        return 'unknown'
    diffs = [b - a for a, b in zip(v, v[1:])]
    if max(diffs) - min(diffs) < 0.15 * (abs(sum(diffs) / len(diffs)) + 1e-9):
        return 'linear'
    if all(a > 0 for a in v):
        ratios = [b / a for a, b in zip(v, v[1:])]
        if max(ratios) - min(ratios) < 0.15 * (sum(ratios) / len(ratios)):
            return 'log'
    return 'irregular'

def extract_graph_structure(text_bounding_boxes):
    """DETERMINISTIC, exact-from-source structure of an X-Y graph: axis titles/units,
    numeric tick VALUES -> range + scale (linear/log), and series/legend labels.
    Read purely from the figure's own text layer via spatial clustering — it never
    traces or invents (x,y) data points, so it cannot fabricate. Returns None if no
    axis structure is confidently recoverable."""
    P = []
    for b in text_bounding_boxes or []:
        t = (b.get('text') or '').replace(' [uncertain]', '').strip()
        if not t:
            continue
        x0, y0, x1, y1 = b['bbox']
        P.append({'t': t, 'x0': x0, 'y0': y0, 'x1': x1, 'y1': y1,
                  'cx': (x0 + x1) / 2, 'cy': (y0 + y1) / 2, 'w': x1 - x0, 'h': y1 - y0})
    if not P:
        return None
    nums = []
    for p in P:
        v = _graph_val(p['t'])
        if v is not None:
            p['v'] = v
            nums.append(p)
    hs = sorted(p['h'] for p in P if p['h'] > 0)
    tol = (hs[len(hs) // 2] * 1.5) if hs else 20.0

    def cluster(items, key):
        clusters = []
        for p in sorted(items, key=lambda p: p[key]):
            for c in clusters:
                if abs(p[key] - c['k']) <= tol:
                    c['items'].append(p)
                    c['k'] = sum(i[key] for i in c['items']) / len(c['items'])
                    break
            else:
                clusters.append({'k': p[key], 'items': [p]})
        return clusters

    def pick_axis(clusters, sortkey):
        for c in sorted(clusters, key=lambda c: -len(c['items'])):
            if len(c['items']) < 3:
                continue
            items = sorted(c['items'], key=lambda p: p[sortkey])
            vals = [p['v'] for p in items]
            if _graph_monotonic(vals):
                return items, {'ticks': vals, 'min': min(vals), 'max': max(vals),
                               'scale': _graph_scale(vals)}
        return None, None

    y_items, y_axis = pick_axis(cluster(nums, 'cx'), 'cy')   # vertical tick column
    x_items, x_axis = pick_axis(cluster(nums, 'cy'), 'cx')   # horizontal tick row

    plot = None
    if y_items:
        plot = (max(p['x1'] for p in y_items), min(p['y0'] for p in y_items),
                max(p['x1'] for p in P), max(p['y1'] for p in y_items))

    texts = [p for p in P if 'v' not in p]
    y_title = x_title = None
    series = []
    for p in texts:
        if y_items and p['h'] > p['w'] * 1.5 and p['x1'] <= min(pp['x0'] for pp in y_items) + tol * 3:
            y_title = p['t']
            continue
        if plot:
            px0, py0, px1, py1 = plot
            if p['cy'] > py1 and p['cx'] >= px0 and len(p['t']) <= 24 and x_title is None:
                x_title = p['t']
                continue
            if p['cx'] >= px0 - tol and py0 - tol <= p['cy'] <= py1 + tol * 4 and len(p['t']) <= 28:
                series.append(p['t'])
    seen, clean = set(), []
    for s in series:
        if s in (x_title, y_title) or s in seen:
            continue
        seen.add(s)
        clean.append(s)
    series = clean

    def unit(title):
        m = re.search(r'\(([^)]+)\)', title or '')
        return m.group(1) if m else None

    if not (y_axis or x_axis):
        return None  # no confident axis structure -> emit nothing rather than guess
    conf = 'high' if (y_axis and (x_axis or x_title)) else 'partial'
    return {
        'x_axis': {**(x_axis or {}), 'title': x_title, 'unit': unit(x_title)},
        'y_axis': {**(y_axis or {}), 'title': y_title, 'unit': unit(y_title)},
        'series_labels': series,
        'confidence': conf,
        'method': 'deterministic-text-layer',
    }

def native_text_boxes_in_rect(page, rect, dpi=300):
    """
    Reads text that is ALREADY digitized in the PDF (exact Unicode glyphs,
    including Greek letters/symbols that OCR frequently mangles) for a given
    region. This must be preferred over OCR for any non-scanned page, since
    the PDF's own text layer is ground truth and OCR is a lossy approximation
    of it. Returns phrases in the same shape cluster_ocr_phrases produces,
    with confidence fixed at 100 (exact, never "[uncertain]").

    Coordinates are converted from PDF points (rect-relative) into the same
    pixel space as the cropped PNG saved via get_pixmap(clip=rect, dpi=dpi),
    since downstream consumers (verify_pixel_connection) index pixels in that
    cropped image directly using these bbox numbers.
    """
    scale = dpi / 72.0
    words_data = page.get_text("words", clip=rect)
    # Drop running headers / page numbers that fall inside a figure crop when the
    # figure sits against a page margin. Done by ABSOLUTE page position (before
    # phrase clustering), so it's robust to the header being split into fragments
    # like "SECTION II" + "• Anesthetic Physiology" — which text matching misses.
    page_h = page.rect.height or 1
    top_margin, bot_margin = page_h * 0.075, page_h * 0.925
    words = []
    for w in words_data:
        text = (w[4] or '').strip()
        if not text:
            continue
        if w[3] < top_margin or w[1] > bot_margin:  # word wholly in top/bottom margin
            continue
        words.append({
            'text': text,
            'x0': (w[0] - rect.x0) * scale,
            'y0': (w[1] - rect.y0) * scale,
            'x1': (w[2] - rect.x0) * scale,
            'y1': (w[3] - rect.y0) * scale,
            'conf': 100.0
        })
    return group_words_into_phrases(words)

def verify_pixel_connection(img_rgb, bbox_a, bbox_b):
    try:
        # Center points of bboxes
        cx_a = int((bbox_a[0] + bbox_a[2]) / 2)
        cy_a = int((bbox_a[1] + bbox_a[3]) / 2)
        cx_b = int((bbox_b[0] + bbox_b[2]) / 2)
        cy_b = int((bbox_b[1] + bbox_b[3]) / 2)
        
        dx = cx_b - cx_a
        dy = cy_b - cy_a
        length = math.sqrt(dx*dx + dy*dy)
        if length < 15:
            return False
            
        # Exclude regions directly inside/near text boxes to avoid text interference
        start_ratio = 15.0 / length
        end_ratio = 1.0 - (15.0 / length)
        if start_ratio >= end_ratio:
            start_ratio = 0.1
            end_ratio = 0.9
            
        total_samples = 40
        width, height = img_rgb.size
        hits = []

        for i in range(total_samples):
            ratio = start_ratio + (end_ratio - start_ratio) * (i / (total_samples - 1))
            px = int(cx_a + dx * ratio)
            py = int(cy_a + dy * ratio)

            px = max(0, min(width - 1, px))
            py = max(0, min(height - 1, py))

            r, g, b = img_rgb.getpixel((px, py))
            # Non-white: dark or saturated pixels
            is_non_white = (r + g + b) / 3.0 < 220 or max(abs(r-g), abs(g-b), abs(r-b)) > 30
            hits.append(is_non_white)

        # A real connector line produces a long, unbroken run of ink along the
        # path. Sparse, scattered hits — the signature of incidental nearby
        # artwork in a busy/dense figure, not an actual line drawn between
        # these two specific labels — are rejected even if they sum to a
        # similar total fraction. This is what previously let figures with
        # lots of surrounding ink (chemical structures, detailed illustrations)
        # produce hundreds of fabricated "connections" between unrelated labels.
        longest_run, current_run = 0, 0
        for h in hits:
            if h:
                current_run += 1
                longest_run = max(longest_run, current_run)
            else:
                current_run = 0

        return (longest_run / total_samples) >= 0.35
    except Exception:
        return False

# ============================================================================
# Vector geometry-based relationship extraction.
#
# Pixel-sampling (verify_pixel_connection above) treats a "line" between two
# labels as evidence of a connection if enough sampled pixels along the
# straight path between them are non-white. That's a coarse proxy: it can't
# tell a real connector line from incidental ink (crossing strokes, nearby
# unrelated glyphs), and it can never know which end has an arrowhead, so it
# always reported a fabricated constant ("Unidirectional"/"Excitatory").
#
# This module instead reads the PDF's own vector path objects directly via
# page.get_drawings() within the figure's crop region:
#   1. Stroke-bearing paths ('l'/'c' items) are flattened into line segments.
#   2. Segment endpoints are clustered (snapped) into junctions, and the
#      resulting graph's connected components are the actual "wires".
#   3. Small filled polygons (2-8 line/curve vertices, small area) sitting at
#      a wire's junction are treated as arrowhead candidates; the polygon
#      vertex farthest from the centroid of the others is its apex, which
#      gives the arrow's pointing direction.
#   4. A wire's loose ends (degree-1 junctions) are matched to the nearest
#      text label; if exactly two (or more, capped) distinct labels are
#      found, that's a real source/target edge, with direction taken from
#      which end (if any) carries an arrowhead — never assumed.
#
# Real-world PDFs vary enormously in how "small filled shape near a line"
# shows up: chemical structure diagrams and detailed illustrations can have
# hundreds to thousands of decorative strokes/fills that would otherwise look
# like a dense web of "connections" if run through this logic uncritically.
# MAX_VECTOR_ELEMENTS exists specifically to detect that case and bail out
# (with a warning) rather than emit relationships for what is actually art,
# not a labeled diagram.
# ============================================================================

WIRE_SNAP_TOL_PX = 5.0       # endpoints within this distance are the same junction
ARROWHEAD_MATCH_TOL_PX = 9.0  # max distance from an arrowhead's centroid to a junction
LABEL_ATTACH_TOL_PX = 14.0    # max distance from a wire's loose end to a text label
PATH_LABEL_TOL_PX = 7.0       # max distance from a label to a segment's midpoint to count as an annotation on that wire
MAX_VECTOR_ELEMENTS = 200     # above this, treat the region as art/chemistry, not a diagram
MAX_LEAVES_PER_COMPONENT = 12  # cap pairwise edges from one wire bundle with many loose ends

def _bezier_polyline(p0, p1, p2, p3, n=4):
    pts = []
    for i in range(n + 1):
        t = i / n
        mt = 1 - t
        x = mt**3 * p0[0] + 3 * mt**2 * t * p1[0] + 3 * mt * t**2 * p2[0] + t**3 * p3[0]
        y = mt**3 * p0[1] + 3 * mt**2 * t * p1[1] + 3 * mt * t**2 * p2[1] + t**3 * p3[1]
        pts.append((x, y))
    return pts

def _to_local(pt, rect, scale):
    return ((pt[0] - rect.x0) * scale, (pt[1] - rect.y0) * scale)

def _flatten_stroke_segments(drawings, rect, scale, max_segments):
    """Flattens stroke-bearing path items within `rect` into local, dpi-scaled
    line segments. Returns None if the region is too complex (more segments
    than max_segments) to plausibly be a labeled diagram rather than art."""
    margin = 2.0
    expanded = fitz.Rect(rect.x0 - margin, rect.y0 - margin, rect.x1 + margin, rect.y1 + margin)
    segments = []
    for d in drawings:
        if d.get('type') not in ('s', 'fs'):
            continue
        if not d['rect'].intersects(expanded):
            continue
        for item in d['items']:
            op = item[0]
            pts = None
            if op == 'l':
                pts = [(item[1].x, item[1].y), (item[2].x, item[2].y)]
            elif op == 'c':
                p0 = (item[1].x, item[1].y)
                p1 = (item[2].x, item[2].y)
                p2 = (item[3].x, item[3].y)
                p3 = (item[4].x, item[4].y)
                pts = _bezier_polyline(p0, p1, p2, p3, n=4)
            if not pts:
                continue
            local_pts = [_to_local(p, rect, scale) for p in pts]
            for i in range(len(local_pts) - 1):
                segments.append((local_pts[i][0], local_pts[i][1], local_pts[i + 1][0], local_pts[i + 1][1]))
                if len(segments) > max_segments:
                    return None
    return segments

def _triangle_fill_candidates(drawings, rect, scale, max_candidates):
    """Finds small filled polygons built from straight/curve segments (NOT
    rectangles — those are almost always data-point markers, not arrowheads)
    within `rect`, and computes each one's apex (the vertex farthest from the
    centroid of the others), which is the direction the arrow points.
    Returns None if there are implausibly many (likely decorative art)."""
    margin = 2.0
    expanded = fitz.Rect(rect.x0 - margin, rect.y0 - margin, rect.x1 + margin, rect.y1 + margin)
    out = []
    for d in drawings:
        if d.get('fill') is None:
            continue
        ops = [it[0] for it in d['items']]
        if not ops or any(o not in ('l', 'c') for o in ops):
            continue
        if not (2 <= len(ops) <= 8):
            continue
        if not d['rect'].intersects(expanded):
            continue
        area = d['rect'].width * d['rect'].height
        if not (0.2 < area < 150):
            continue

        verts = []
        for it in d['items']:
            if it[0] == 'l':
                verts.append((it[1].x, it[1].y))
                verts.append((it[2].x, it[2].y))
            elif it[0] == 'c':
                verts.append((it[1].x, it[1].y))
                verts.append((it[4].x, it[4].y))

        uniq = []
        for v in verts:
            if not uniq or math.hypot(v[0] - uniq[-1][0], v[1] - uniq[-1][1]) > 0.05:
                uniq.append(v)
        if len(uniq) < 3:
            continue

        cx = sum(v[0] for v in uniq) / len(uniq)
        cy = sum(v[1] for v in uniq) / len(uniq)
        apex, best_score = uniq[0], -1.0
        for i, v in enumerate(uniq):
            others = uniq[:i] + uniq[i + 1:]
            bx = sum(o[0] for o in others) / len(others)
            by = sum(o[1] for o in others) / len(others)
            score = math.hypot(v[0] - bx, v[1] - by)
            if score > best_score:
                best_score = score
                apex = v

        out.append({
            'centroid': _to_local((cx, cy), rect, scale),
            'apex': _to_local(apex, rect, scale)
        })
        if len(out) > max_candidates:
            return None
    return out

def _cluster_endpoints(points, tol):
    """Union-find clustering of nearby points (segment endpoints) into
    junctions, so a connector drawn as several short line/curve pieces is
    treated as one continuous wire instead of disjoint fragments."""
    n = len(points)
    parent = list(range(n))

    def find(i):
        while parent[i] != i:
            parent[i] = parent[parent[i]]
            i = parent[i]
        return i

    def union(i, j):
        ri, rj = find(i), find(j)
        if ri != rj:
            parent[ri] = rj

    for i in range(n):
        xi, yi = points[i]
        for j in range(i + 1, n):
            xj, yj = points[j]
            if abs(xi - xj) <= tol and abs(yi - yj) <= tol and math.hypot(xi - xj, yi - yj) <= tol:
                union(i, j)

    roots = [find(i) for i in range(n)]
    remap = {}
    cluster_ids = []
    for r in roots:
        if r not in remap:
            remap[r] = len(remap)
        cluster_ids.append(remap[r])
    return cluster_ids, len(remap)

def _is_substantive_label(text):
    """
    A meaningful node/connection label must contain at least one letter.
    Without this, a chart's axis line or gridline — which is real vector
    geometry connecting two real text boxes at its ends — gets reported as a
    "relationship" between two axis tick numbers (e.g. "500" -> "-10"). That's
    a artifact of chart decoration, not a labeled connection between entities.
    """
    return bool(re.search(r'[A-Za-zΑ-Ωα-ω]', text))

def _label_for_point(point, text_bounding_boxes, tol):
    px, py = point
    best_label, best_dist = None, float('inf')
    for b in text_bounding_boxes:
        bx0, by0, bx1, by1 = b['bbox']
        dx = max(bx0 - px, 0.0, px - bx1)
        dy = max(by0 - py, 0.0, py - by1)
        dist = math.hypot(dx, dy)
        if dist < best_dist:
            best_dist = dist
            best_label = b
    if best_label is not None and best_dist <= tol:
        return best_label['text'].replace(" [uncertain]", "")
    return None

def extract_vector_relationships(drawings, rect, scale, text_bounding_boxes, warnings, fig_id, page_num):
    """
    Real geometry-based replacement for the pixel-sampling heuristic: traces
    actual PDF vector paths within a figure's crop region to find connector
    wires between labeled nodes, and determines real direction from detected
    arrowhead polygons rather than asserting a constant. Returns [] (with a
    logged warning) if the region is too vector-dense to plausibly be a
    labeled diagram, and [] silently if there's simply nothing to find.
    """
    if not text_bounding_boxes or not drawings:
        return []

    segments = _flatten_stroke_segments(drawings, rect, scale, MAX_VECTOR_ELEMENTS)
    if segments is None:
        warnings.append(
            f"Figure {fig_id} on page {page_num} has more than {MAX_VECTOR_ELEMENTS} "
            "vector line segments in its region (likely a chemical structure or "
            "detailed illustration, not a labeled diagram) — relationship tracing was skipped."
        )
        return []
    if not segments:
        return []

    fills = _triangle_fill_candidates(drawings, rect, scale, MAX_VECTOR_ELEMENTS)
    if fills is None:
        warnings.append(
            f"Figure {fig_id} on page {page_num} has more than {MAX_VECTOR_ELEMENTS} "
            "small filled vector shapes in its region — relationship tracing was skipped."
        )
        return []

    endpoints = []
    for (x0, y0, x1, y1) in segments:
        endpoints.append((x0, y0))
        endpoints.append((x1, y1))

    cluster_ids, n_clusters = _cluster_endpoints(endpoints, WIRE_SNAP_TOL_PX)

    cluster_adj = {c: set() for c in range(n_clusters)}
    seg_clusters = []
    for si in range(len(segments)):
        ca, cb = cluster_ids[si * 2], cluster_ids[si * 2 + 1]
        seg_clusters.append((ca, cb))
        if ca != cb:
            cluster_adj[ca].add(cb)
            cluster_adj[cb].add(ca)

    sums = {}
    counts = {}
    for idx, pt in enumerate(endpoints):
        cid = cluster_ids[idx]
        sx, sy = sums.get(cid, (0.0, 0.0))
        sums[cid] = (sx + pt[0], sy + pt[1])
        counts[cid] = counts.get(cid, 0) + 1
    cluster_centroid = {cid: (sx / counts[cid], sy / counts[cid]) for cid, (sx, sy) in sums.items()}
    cluster_degree = {c: len(cluster_adj[c]) for c in range(n_clusters)}

    # Match arrowhead fill candidates to the nearest junction within tolerance.
    head_clusters = set()
    for f in fills:
        best_c, best_d = None, float('inf')
        fx, fy = f['centroid']
        for cid, (cx, cy) in cluster_centroid.items():
            d = math.hypot(cx - fx, cy - fy)
            if d < best_d:
                best_d = d
                best_c = cid
        if best_c is not None and best_d <= ARROWHEAD_MATCH_TOL_PX:
            head_clusters.add(best_c)

    # Connected components over the junction graph = logical wires.
    visited = set()
    components = []
    for c in range(n_clusters):
        if c in visited:
            continue
        stack, comp = [c], set()
        while stack:
            cur = stack.pop()
            if cur in comp:
                continue
            comp.add(cur)
            visited.add(cur)
            stack.extend(nb for nb in cluster_adj[cur] if nb not in comp)
        components.append(comp)

    edges = []
    seen_pairs = set()
    for comp in components:
        leaves = [c for c in comp if cluster_degree.get(c, 0) <= 1]
        leaf_labels = []
        for c in leaves:
            label = _label_for_point(cluster_centroid[c], text_bounding_boxes, LABEL_ATTACH_TOL_PX)
            if label and _is_substantive_label(label):
                leaf_labels.append((c, label))
        if len(leaf_labels) < 2:
            continue
        leaf_labels = leaf_labels[:MAX_LEAVES_PER_COMPONENT]

        # Text sitting on/near any segment in this wire (not just its ends) is
        # treated as an annotation of the connection itself (e.g. "GABA (-)"
        # written along the line), folded into sign/neurotransmitter inference.
        path_label_texts = []
        for si, (ca, cb) in enumerate(seg_clusters):
            if ca not in comp and cb not in comp:
                continue
            x0, y0, x1, y1 = segments[si]
            mid = ((x0 + x1) / 2, (y0 + y1) / 2)
            label = _label_for_point(mid, text_bounding_boxes, PATH_LABEL_TOL_PX)
            if label:
                path_label_texts.append(label)
        path_text = ' '.join(path_label_texts)

        for i in range(len(leaf_labels)):
            for j in range(i + 1, len(leaf_labels)):
                ca, label_a = leaf_labels[i]
                cb, label_b = leaf_labels[j]
                if label_a == label_b:
                    continue
                pair_key = tuple(sorted([label_a, label_b]))
                if pair_key in seen_pairs:
                    continue
                seen_pairs.add(pair_key)

                head_a, head_b = ca in head_clusters, cb in head_clusters
                if head_a and head_b:
                    source, target, direction = label_a, label_b, "Bidirectional"
                elif head_b and not head_a:
                    source, target, direction = label_a, label_b, "Unidirectional"
                elif head_a and not head_b:
                    source, target, direction = label_b, label_a, "Unidirectional"
                else:
                    source, target, direction = label_a, label_b, "Undirected"

                edges.append({
                    "source": source,
                    "target": target,
                    "direction": direction,
                    "status": "vector_geometry_verified",
                    **infer_sign_and_neurotransmitter(label_a, label_b, path_text)
                })

    return edges

def extract_figures_from_scanned_page(page, page_num, source_file, scratch_dir, word_boxes, page_img, scale_x, scale_y):
    """
    Best-effort figure capture for scanned/raster pages. Unlike born-digital
    pages, there's no embedded-image or vector-drawing metadata to anchor on
    (the whole page IS one image), so this groups OCR'd words into lines,
    finds caption-like lines ("Fig. 9-1", "Table 9-1"), and crops the nearest
    large whitespace gap adjacent to that caption as the figure region.
    Deliberately conservative: skips a caption rather than guessing a region
    when no clear gap exists. This is necessarily more limited than the
    born-digital path (no vector relationship tracing is attempted here).
    """
    if not word_boxes:
        return []

    sorted_words = sorted(word_boxes, key=lambda w: (w['bbox'][1], w['bbox'][0]))
    lines = []
    for w in sorted_words:
        placed = False
        for line in lines:
            if abs(w['bbox'][1] - line['y0']) < 8:
                line['words'].append(w)
                line['y1'] = max(line['y1'], w['bbox'][3])
                placed = True
                break
        if not placed:
            lines.append({'y0': w['bbox'][1], 'y1': w['bbox'][3], 'words': [w]})
    for line in lines:
        line['words'].sort(key=lambda w: w['bbox'][0])
        line['text'] = ' '.join(w['text'] for w in line['words'])
        line['x0'] = min(w['bbox'][0] for w in line['words'])
        line['x1'] = max(w['bbox'][2] for w in line['words'])
    lines.sort(key=lambda l: l['y0'])

    caption_idxs = [i for i, l in enumerate(lines) if re.match(r'^(?:Fig\b\.|Figure\b|Table\b)', l['text'], re.IGNORECASE)]

    figures = []
    page_h, page_w = page.rect.height, page.rect.width
    mid_x = page_w / 2.0

    for c_idx, li in enumerate(caption_idxs):
        cap_line = lines[li]
        cap_text = cap_line['text']
        column_left = (cap_line['x0'] + cap_line['x1']) / 2 < mid_x
        col_x0, col_x1 = (0.0, mid_x) if column_left else (mid_x, page_w)

        def _col_center(l):
            return (l['x0'] + l['x1']) / 2

        prev_bound = max(
            (lines[j]['y1'] for j in range(li - 1, -1, -1) if col_x0 <= _col_center(lines[j]) <= col_x1),
            default=0.0
        )
        next_bound = min(
            (lines[j]['y0'] for j in range(li + 1, len(lines)) if col_x0 <= _col_center(lines[j]) <= col_x1),
            default=page_h
        )

        gap_above = cap_line['y0'] - prev_bound
        gap_below = next_bound - cap_line['y1']

        if max(gap_above, gap_below) < 40:
            continue  # No meaningful whitespace gap near this caption — skip rather than guess.

        if gap_above >= gap_below:
            y0, y1 = max(0.0, cap_line['y0'] - min(gap_above, 400)), cap_line['y0'] - 2
        else:
            y0, y1 = cap_line['y1'] + 2, min(page_h, cap_line['y1'] + min(gap_below, 400))

        rect_pt = fitz.Rect(col_x0 + 5, y0, col_x1 - 5, y1)
        if rect_pt.height < 20 or rect_pt.width < 20:
            continue

        px = (int(rect_pt.x0 / scale_x), int(rect_pt.y0 / scale_y),
              int(rect_pt.x1 / scale_x), int(rect_pt.y1 / scale_y))
        try:
            crop_img = page_img.crop(px)
        except Exception:
            continue

        fig_id = sanitize_fig_id(cap_text, f"scan_{c_idx}", page_num)
        out_filename = f"{source_file.replace('.', '_')}_{fig_id}.png"
        out_path = os.path.join(scratch_dir, out_filename)
        try:
            crop_img.save(out_path)
        except Exception:
            continue

        # No PDF text layer exists on a scanned page, so OCR is the only option here.
        text_bounding_boxes = strip_furniture_phrases(cluster_ocr_phrases(out_path))
        cap_lower = cap_text.lower()
        archetype = "COORDINATE X-Y GRAPHS & COMPLEMENTARY PANELS"
        # This initial guess gets superseded by the TS-side StrategyRouter's own
        # (more complete) keyword match — see archetypes.ts's CONTINUOUS_WAVEFORM
        # config — but keep it modality-neutral so it isn't misleading in the interim.
        if any(kw in cap_lower for kw in (
            "eeg", "electroencephalogram", "ecg", "ekg", "electrocardiogram",
            "capnogr", "etco2", "plethysmogr", "waveform", "tracing"
        )):
            archetype = "CONTINUOUS_WAVEFORM"
        elif "hypnogram" in cap_lower:
            archetype = "TIMELINE_STEP_CHART_HYPNOGRAM"

        figures.append({
            "id": fig_id,
            "sourceFile": source_file,
            "pageNumber": page_num,
            "image_path": os.path.abspath(out_path),
            "closest_text_heading": "",
            "caption": cap_text,
            "text_bounding_boxes": text_bounding_boxes,
            "archetype": archetype,
            "details": {
                "labels": [b["text"].replace(" [uncertain]", "") for b in text_bounding_boxes]
            }
        })

    return figures

def find_gutter_x(blocks, page_w):
    """For a two-column page, return the x of the vertical gutter, adaptively (so
    asymmetric wide+narrow columns work), or None if the page is single-column.

    Looks only in the central band and only at column-width prose blocks (not
    full-width headers/figures): the gutter is a near-empty vertical strip there.
    Returning None → caller reads the page as a single column. This is deliberately
    conservative: it enhances the proven center split, it does not replace it."""
    bins = 80
    binw = page_w / bins
    cover = [0] * bins
    for b in blocks:
        if (b[2] - b[0]) > page_w * 0.6:  # skip spanning/full-width blocks
            continue
        i0 = max(0, int(b[0] / binw))
        i1 = min(bins - 1, int(b[2] / binw))
        for i in range(i0, i1 + 1):
            cover[i] += 1
    lo_i, hi_i = int(bins * 0.38), int(bins * 0.62)
    central = [(i, cover[i]) for i in range(lo_i, hi_i + 1)]
    if not central or max(c for _, c in central) == 0:
        return None
    min_cov = min(c for _, c in central)
    side_max = max(cover[:lo_i] + cover[hi_i + 1:]) if bins > (hi_i + 1) else 0
    # A real gutter: the central strip is (near-)empty while the sides have text.
    if min_cov > 0 or side_max < 2:
        return None
    empties = [i for i, c in central if c == min_cov]
    return (min(empties) + max(empties) + 1) / 2.0 * binw

def order_blocks_reading_order(blocks, page_w):
    """Return blocks in true human reading order for 1- or 2-column pages.

    A plain (y, x) sort interleaves the columns of a two-column page: it emits a
    line from the left column, then the line beside it in the right column, and so
    on — splicing unrelated half-sentences together. Correct order reads the ENTIRE
    left column top-to-bottom, then the right column. Full-width "spanning" blocks
    (titles, cross-column figures/tables) split the page into horizontal bands; each
    band is read left column then right. Single-column pages fall through to plain
    top-to-bottom. The gutter x is detected adaptively (find_gutter_x), defaulting to
    page center — reliable for the 1- and 2-column layouts textbooks actually use.
    (True 3+ column support is intentionally deferred until a 3-column source exists
    to verify against — see docs/parser_hardening_audit.md item 2.2.)
    """
    if not blocks or page_w <= 0:
        return blocks
    mid = find_gutter_x(blocks, page_w)
    if mid is None:
        mid = page_w / 2.0
    gutter_tol = page_w * 0.02
    spanning, left, right = [], [], []
    for b in blocks:
        x0, x1 = b[0], b[2]
        width = x1 - x0
        center = (x0 + x1) / 2.0
        crosses_gutter = x0 < mid - gutter_tol and x1 > mid + gutter_tol
        if crosses_gutter and width > page_w * 0.55:
            spanning.append(b)
        elif center < mid:
            left.append(b)
        else:
            right.append(b)

    # No real two-column structure -> keep simple top-to-bottom, left-to-right.
    if not left or not right:
        return sorted(blocks, key=lambda b: (b[1], b[0]))

    spanning.sort(key=lambda b: b[1])
    sep_ys = [b[1] for b in spanning]
    band_of = lambda b: sum(1 for sy in sep_ys if sy <= b[1])

    from collections import defaultdict
    bands_left, bands_right = defaultdict(list), defaultdict(list)
    for b in left:
        bands_left[band_of(b)].append(b)
    for b in right:
        bands_right[band_of(b)].append(b)

    ordered = []
    for band in range(len(spanning) + 1):
        ordered.extend(sorted(bands_left.get(band, []), key=lambda b: b[1]))
        ordered.extend(sorted(bands_right.get(band, []), key=lambda b: b[1]))
        if band < len(spanning):
            ordered.append(spanning[band])
    return ordered

# Common English function/domain words used only as a cheap "is this real English
# prose or CID-garbled/non-English gibberish?" signal (never to alter text).
_COMMON_WORDS = set((
    "the of and to in a is that for it as was with be by on not he this are or from at "
    "which but have an they you one had who all will more no if out so up what its about "
    "into than them can only other new some could time these two may then do first any my "
    "now such like our over me even most made after also did many before must through back "
    "where much your way down should because each just those people how between both under "
    "result increase decrease patients blood pressure effects during anesthesia dose cause "
    "use administration heart rate cardiac renal hepatic respiratory arterial venous"
).split())

# Ligatures some fonts encode as single codepoints -> ASCII equivalents.
_LIGATURES = {0xFB00: "ff", 0xFB01: "fi", 0xFB02: "fl", 0xFB03: "ffi",
              0xFB04: "ffl", 0xFB05: "ft", 0xFB06: "st"}

# Figure-caption openers across common textbook styles + a few non-English forms,
# so figure capture isn't limited to English "Fig./Figure".
_FIG_CAPTION_RE = re.compile(
    r"^(?:Fig(?:ure)?\b\.?|Abb\b\.?|Figura\b|Scheme\b|Plate\b|Chart\b|Box\b|Exhibit\b)",
    re.IGNORECASE)

_TABLE_CAPTION_RE = re.compile(r"\bTab(?:le|\.)\s*\d", re.IGNORECASE)

def page_has_table_signal(raw_blocks, drawings, page_w):
    """Is a real table plausible on this page? Real data tables carry a 'Table N'
    caption or ruling lines; prose/figure pages have neither. `find_tables()` is by
    far the most expensive per-page call (~75% of parse time) and it also invents
    false 'tables' from dense 2-column prose — so we only run it where this cheap
    predicate says a table is plausible. Kept deliberately permissive (2 ruling
    marks) so borderline real tables are not skipped."""
    for b in raw_blocks:
        if b[6] == 0 and _TABLE_CAPTION_RE.search(b[4] or ""):
            return True
    ruled = 0
    for d in drawings or []:
        for it in d.get("items", []):
            if it[0] == "l":
                p1, p2 = it[1], it[2]
                if abs(p1.y - p2.y) < 1.5 and abs(p1.x - p2.x) > page_w * 0.2:
                    ruled += 1
            elif it[0] == "re":
                r = it[1]
                if r.width > page_w * 0.2 and r.height > 4:
                    ruled += 1
            if ruled >= 2:
                return True
    return False

def clean_text(s):
    """Normalize extracted text for fidelity, WITHOUT reordering or dropping words:
      - Unicode NFC (canonical form) + expand ﬁ/ﬂ-style ligatures,
      - repair soft line-break hyphenation ("anes-\\nthetics" -> "anesthetics").
    De-hyphenation only fires on lowercase-hyphen-newline-lowercase, the syllable-
    break pattern; a hyphen followed by a capital or digit (real compound / range)
    is left intact."""
    if not s:
        return s
    s = unicodedata.normalize("NFC", s).translate(_LIGATURES)
    s = re.sub(r"([a-z])-\n[ \t]*([a-z])", r"\1\2", s)
    return s

# A superscript that is clearly a citation/footnote marker: 2+ digits, or
# comma/dash-joined digit groups ("28,29"). Single digits are left alone so
# exponents/units ("cm2", "m2") and chemistry ("CO2") are never broken.
_SUPERSCRIPT_CITATION_RE = re.compile(r"^\d{2,}$|^\d+(?:[,–-]\d+)+$")

def superscript_separated_block_texts(page):
    """Map block-number -> text with citation superscripts separated from the word
    they were glued to ('movement28,29' -> 'movement 28,29'), so tokens/search stay
    clean. Uses the PyMuPDF superscript span flag; only multi-digit/comma markers are
    touched. Returns {} on any failure (caller then keeps the original block text)."""
    out = {}
    try:
        dd = page.get_text("dict")
    except Exception:
        return out
    for b in dd.get("blocks", []):
        if b.get("type", 0) != 0:
            continue
        line_texts = []
        for l in b.get("lines", []):
            buf = ""
            for s in l.get("spans", []):
                t = s.get("text", "")
                if (s.get("flags", 0) & 1) and buf and buf[-1] not in " \t" \
                        and _SUPERSCRIPT_CITATION_RE.match(t.strip()):
                    buf += " " + t
                else:
                    buf += t
            line_texts.append(buf)
        out[b.get("number")] = "\n".join(line_texts)
    return out

def _furniture_sig(s):
    """Normalize a candidate header/footer line so page numbers vary but the
    template matches: 'SECTION II • Physiology 491' -> 'sectioniiphysiology'."""
    return re.sub(r"[^a-z]", "", s.lower())

def detect_running_furniture(doc):
    """Find running headers/footers: text lines that recur near the top or bottom
    margin on many pages (page-number digits ignored), so this print furniture never
    pollutes the prose or the knowledge DB.

    Returns (furniture_signatures, blocks_by_page): the raw get_text('blocks') result
    for every page is cached and returned so the main pass reuses it instead of
    re-extracting (blocks were otherwise extracted twice per page)."""
    from collections import Counter
    top, bot = Counter(), Counter()
    npages = len(doc)
    blocks_by_page = []
    for page in doc:
        h = page.rect.height or 1
        try:
            raw = page.get_text("blocks")
        except Exception:
            raw = []
        blocks_by_page.append(raw)
        for b in raw:
            if b[6] != 0 or not b[4].strip():
                continue
            first = b[4].strip().split("\n")[0].strip()
            sig = _furniture_sig(first)
            if len(sig) < 3:
                continue
            if b[3] < h * 0.12:
                top[sig] += 1
            elif b[1] > h * 0.88:
                bot[sig] += 1
    thresh = max(3, int(npages * 0.3))
    furniture = {s for s, c in top.items() if c >= thresh} | {s for s, c in bot.items() if c >= thresh}
    return furniture, blocks_by_page

def is_furniture_block(b, page_h, furniture):
    """A block is print furniture if it's a bare page number, or a known running
    header/footer signature sitting in the top/bottom margin."""
    text = b[4].strip()
    if re.fullmatch(r"\d{1,4}", text):
        return True
    sig = _furniture_sig(text.split("\n")[0].strip())
    if sig in furniture and (b[3] < page_h * 0.12 or b[1] > page_h * 0.88):
        return True
    return False

def page_needs_ocr(page):
    """Per-PAGE decision (not per-document): OCR a page only when its native text
    layer is empty/sparse OR is CID-garbled gibberish. This handles hybrid PDFs
    (some real-text pages, some scanned) that a single per-document flag mishandles
    — and recovers pages whose fonts have no usable ToUnicode map."""
    try:
        txt = page.get_text().strip()
    except Exception:
        return True
    if len(txt) < 30:
        return True  # effectively no text layer on this page
    toks = re.findall(r"[A-Za-z]+", txt)
    if len(toks) > 40:
        eng = sum(1 for t in toks if t.lower() in _COMMON_WORDS) / len(toks)
        if eng < 0.08:
            return True  # decodes to gibberish -> OCR the rendered image instead
    return False

def compute_quality_metrics(fragments):
    """Cheap, deterministic text-quality signals so a bad parse is never silent.
    Operates on the already-extracted text; never changes it."""
    text = "\n".join(f.get("rawText", "") or "" for f in fragments)
    n = len(text) or 1
    pages = len(fragments) or 1
    total_chars = sum(f.get("characterCount", 0) for f in fragments)
    empty = sum(1 for f in fragments if f.get("characterCount", 0) == 0)
    low = sum(1 for f in fragments if 0 < f.get("characterCount", 0) < 150)
    nonascii = sum(1 for c in text if ord(c) > 127)
    repl = text.count("�")
    toks = re.findall(r"[A-Za-z]+", text)
    low_toks = [t.lower() for t in toks]
    eng = (sum(1 for t in low_toks if t in _COMMON_WORDS) / len(low_toks)) if low_toks else 0.0
    longtok = (sum(1 for t in toks if len(t) > 22) / len(toks)) if toks else 0.0
    cpp = total_chars / pages
    flags = []
    if total_chars == 0:
        flags.append("NO_TEXT")
    elif cpp < 400:
        flags.append("LOW_TEXT")
    if eng < 0.15 and len(toks) > 200:
        flags.append("GARBLED_OR_NON_ENGLISH")
    if nonascii / n > 0.08:
        flags.append("HIGH_NONASCII")
    if repl > 0:
        flags.append("DECODE_ERRORS")
    if longtok > 0.03:
        flags.append("RUN_TOGETHER_WORDS")
    if empty > pages * 0.25:
        flags.append("MANY_EMPTY_PAGES")
    critical = {"NO_TEXT", "GARBLED_OR_NON_ENGLISH", "DECODE_ERRORS"}
    return {
        "pages": pages,
        "total_characters": total_chars,
        "chars_per_page": round(cpp, 1),
        "empty_pages": empty,
        "low_text_pages": low,
        "english_word_ratio": round(eng, 3),
        "nonascii_pct": round(100 * nonascii / n, 2),
        "replacement_chars": repl,
        "run_together_pct": round(100 * longtok, 2),
        "flags": flags,
        "ok": not any(f in critical for f in flags),
    }

def extract_pdf(file_path):
    doc = fitz.open(file_path)
    source_file = os.path.basename(file_path)

    # Encrypted/permission-locked PDFs: an empty password unlocks the common
    # "owner-permissions only" case; a real user password can't be recovered, so
    # warn explicitly instead of silently returning empty text.
    enc_warning = None
    if getattr(doc, "is_encrypted", False):
        try:
            authed = doc.authenticate("")
        except Exception:
            authed = 0
        if not authed:
            enc_warning = ("PDF is encrypted/password-protected and could not be opened without a "
                           "password — no text can be extracted. Provide a decrypted copy.")

    # Whether a page uses OCR is decided PER PAGE below (see page_needs_ocr),
    # so hybrid PDFs (mixed real-text and scanned pages) are handled correctly.
    running_furniture, blocks_by_page = detect_running_furniture(doc)

    fragments = []
    visual_data_engines = []
    warnings = []
    if enc_warning:
        warnings.append(enc_warning)
    
    try:
        if doc.is_repaired:
            warnings.append("The PDF cross-reference table was corrupt and was automatically repaired.")
    except AttributeError:
        pass
        
    scratch_dir = SCRATCH_DIR
    os.makedirs(scratch_dir, exist_ok=True)

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_num = page_idx + 1
        page_id = f"PAGE_{page_num:03d}"
        
        raw_text = ""
        sections = []
        word_bounding_boxes = []

        if page_needs_ocr(page):
            try:
                pix = page.get_pixmap(dpi=300)
                img_data = pix.tobytes("png")
                from io import BytesIO
                img = Image.open(BytesIO(img_data))
                
                # OCR language is configurable for non-English books (OCR_LANG,
                # e.g. "deu" or "eng+deu"); psm 3 = automatic page segmentation so
                # Tesseract detects columns/layout instead of reading straight across.
                lang = os.environ.get('OCR_LANG') or (
                    'eng_best' if os.path.exists('/opt/homebrew/share/tessdata/eng_best.traineddata') else 'eng')
                custom_config = f'--oem 1 --psm 3 -l {lang}'

                # Word coordinates via OCR
                data = pytesseract.image_to_data(img, config=custom_config, output_type=pytesseract.Output.DICT)
                raw_text_parts = []
                scale_x = page.rect.width / img.width
                scale_y = page.rect.height / img.height
                
                for i in range(len(data['text'])):
                    word = data['text'][i].strip()
                    try:
                        c = float(data['conf'][i])
                    except (ValueError, TypeError):
                        c = 0.0
                    if word and c > 0:
                        x0 = round(data['left'][i] * scale_x, 1)
                        y0 = round(data['top'][i] * scale_y, 1)
                        x1 = round((data['left'][i] + data['width'][i]) * scale_x, 1)
                        y1 = round((data['top'][i] + data['height'][i]) * scale_y, 1)
                        
                        word_bounding_boxes.append({
                            "text": word,
                            "bbox": [x0, y0, x1, y1]
                        })
                        raw_text_parts.append(word)
                raw_text = clean_text(" ".join(raw_text_parts))

                try:
                    scanned_figs = extract_figures_from_scanned_page(
                        page, page_num, source_file, scratch_dir, word_bounding_boxes, img, scale_x, scale_y
                    )
                    visual_data_engines.extend(scanned_figs)
                except Exception as fig_err:
                    warnings.append(f"Scanned-page figure capture failed on page {page_num}: {str(fig_err)}")
            except Exception as e:
                warnings.append(f"OCR failed on Page {page_num}: {str(e)}")
                raw_text = page.get_text() or ""
        else:
            # Word coordinates via PyMuPDF
            words_data = page.get_text("words")
            for w in words_data:
                word_bounding_boxes.append({
                    "text": w[4],
                    "bbox": [round(w[0], 1), round(w[1], 1), round(w[2], 1), round(w[3], 1)]
                })

            # Reuse blocks already extracted during the furniture pre-pass.
            raw_blocks = sorted(blocks_by_page[page_idx], key=lambda b: (b[1], b[0]))

            # Drawings computed once here (reused by figure extraction below via the
            # _page_drawings_cache) and used to decide whether the costly find_tables()
            # is even worth running on this page.
            page_drawings = page.get_drawings()

            # --- Native table detection (deterministic: vector ruling lines /
            # whitespace gaps, no OCR involved) runs BEFORE the flowing-text pass,
            # so grid structure is captured as real rows/columns instead of being
            # flattened into linear reading-order prose and lost permanently.
            #
            # find_tables() is prone to false positives on a dense 2-column medical
            # layout: "Key Points" sidebar boxes and figure/diagram label clusters
            # both have ruling-line-like geometry that gets mis-detected as a table.
            # Empirically (verified against real chapter pages), every false positive
            # shares one of two signatures: (a) mostly-empty cells (a real data table
            # is densely filled), or (b) one giant cell holding paragraph/bullet text
            # (a real table cell holds a short value, not a paragraph). Both gates
            # below exist specifically to reject those, plus excluding any candidate
            # that overlaps an embedded image (diagrams, not tables, live there). ---
            images_for_table_filter = page.get_image_info(xrefs=True)
            table_rects = []
            # Only run the expensive find_tables() where a table is actually plausible.
            try:
                table_finder = (page.find_tables()
                                if page_has_table_signal(raw_blocks, page_drawings, page.rect.width)
                                else None)
                for t_idx, table in enumerate(table_finder.tables if table_finder else []):
                    try:
                        extracted_rows = table.extract()
                    except Exception:
                        extracted_rows = []
                    if len(extracted_rows) < 2:
                        continue

                    # Clean cell text (NFC/ligatures/de-hyphenation) for table fidelity.
                    rows = [[clean_text(cell or '').strip() for cell in row] for row in extracted_rows]
                    headers = rows[0]

                    total_cells = sum(len(r) for r in rows)
                    nonempty_cells = sum(1 for r in rows for c in r if c)
                    fill_ratio = (nonempty_cells / total_cells) if total_cells else 0.0
                    max_cell_len = max((len(c) for r in rows for c in r), default=0)
                    t_rect_check = fitz.Rect(table.bbox)
                    overlaps_image = any(
                        rect_overlap_ratio(t_rect_check, fitz.Rect(img['bbox'])) > 0.3
                        for img in images_for_table_filter
                    )
                    if fill_ratio < 0.35 or max_cell_len > 200 or overlaps_image:
                        continue
                    # A stray "|" or newline inside a cell would corrupt the Markdown
                    # table (shifting/merging columns downstream) — make cells safe.
                    def _md_cell(c):
                        return c.replace("\n", " ").replace("|", "\\|").strip()
                    markdown_table = "| " + " | ".join(_md_cell(h) for h in headers) + " |\n"
                    markdown_table += "| " + " | ".join(["---"] * len(headers)) + " |\n"
                    for row in rows[1:]:
                        padded = list(row) + [''] * (len(headers) - len(row))
                        markdown_table += "| " + " | ".join(_md_cell(x) for x in padded[:len(headers)]) + " |\n"

                    t_rect = fitz.Rect(table.bbox)
                    table_rects.append(t_rect)

                    caption_text = find_nearby_caption(raw_blocks, t_rect, r'^Table\b')
                    table_id = f"TBL_PAGE_{page_num:03d}_{t_idx}"

                    visual_data_engines.append({
                        "id": table_id,
                        "sourceFile": source_file,
                        "pageNumber": page_num,
                        "image_path": "",
                        "closest_text_heading": closest_heading_before_y(raw_blocks, t_rect.y0),
                        "caption": caption_text or f"Table on page {page_num}",
                        "text_bounding_boxes": [],
                        "archetype": "COORDINATE X-Y GRAPHS & COMPLEMENTARY PANELS",
                        "details": {
                            "matrix_rows": rows,
                            "markdown_representation": markdown_table,
                            "headers": headers,
                            "labels": [c for row in rows for c in row if c]
                        }
                    })
            except Exception as e:
                warnings.append(f"Table detection failed on page {page_num}: {str(e)}")

            # Exclude any flowing-text block that's mostly inside a detected table —
            # otherwise the same content appears twice: once correctly structured
            # above, once flattened/garbled into linear paragraph text below.
            _page_h = page.rect.height
            blocks = [
                b for b in raw_blocks
                if not (b[6] == 0 and any(rect_overlap_ratio(fitz.Rect(b[:4]), tr) > 0.5 for tr in table_rects))
                and not (b[6] == 0 and is_furniture_block(b, _page_h, running_furniture))
            ]

            # Re-order into true reading order (whole left column, then whole
            # right column) so a two-column page is not interleaved. The raw
            # (y, x) sort above is kept for table geometry; text/sections below
            # must follow human reading order instead.
            blocks = order_blocks_reading_order(blocks, page.rect.width)

            # Per-block text with citation superscripts separated (keyed by block
            # number); .get() falls back to the original text if unavailable.
            super_texts = superscript_separated_block_texts(page)
            block_text = lambda b: clean_text(super_texts.get(b[5], b[4]))

            text_parts = []
            captions = []

            for idx, block in enumerate(blocks):
                if block[6] == 0:
                    text = block_text(block).strip()
                    if not text:
                        continue
                    text_parts.append(text)
                    if _FIG_CAPTION_RE.match(text):
                        captions.append((idx, block, text))

            raw_text = "\n\n".join(text_parts)
            
            # Heading matching
            current_sec = None
            for idx, block in enumerate(blocks):
                if block[6] == 0:
                    text = block_text(block).strip()
                    lines = [l.strip() for l in text.split('\n') if l.strip()]
                    if lines:
                        first_line = lines[0]
                        is_heading = is_real_heading(first_line)

                        if is_heading:
                            if current_sec:
                                sections.append(current_sec)
                            current_sec = {
                                "heading": first_line,
                                "body": "\n".join(lines[1:]) if len(lines) > 1 else "",
                                "startLine": len(sections) + 1,
                                "category": "heading"
                            }
                        else:
                            if current_sec:
                                current_sec["body"] += "\n" + text
                            else:
                                current_sec = {
                                    "heading": "",
                                    "body": text,
                                    "startLine": len(sections) + 1,
                                    "category": "paragraph"
                                }
            if current_sec:
                sections.append(current_sec)
                
            # Perform programmatic visual cropping (Phase 1)
            images = images_for_table_filter
            mid_x = page.rect.width / 2.0

            # Reuse the drawings already fetched for the table-signal check above.
            _page_drawings_cache = [page_drawings]

            def get_page_drawings():
                if not _page_drawings_cache:
                    _page_drawings_cache.append(page.get_drawings())
                return _page_drawings_cache[0]

            def _column_of(bbox):
                cx = (bbox[0] + bbox[2]) / 2.0
                return 'left' if cx < mid_x else 'right'

            # Greedy nearest-neighbor assignment of embedded images to captions,
            # computed globally across the whole page BEFORE any cropping happens.
            # Matching each caption independently against "whichever image is
            # closest" lets two different captions both claim the SAME image when
            # they sit close together on the page (e.g. two stacked figures sharing
            # one nearby embedded diagram) — exactly the bug that caused two
            # genuinely distinct figures (an ECG strip and a lead-placement diagram)
            # to crop to identical pixels. Each image can satisfy at most one
            # caption; the globally-closest (caption, image) pairing wins first.
            all_candidate_pairs = []
            for cap_idx, cap_block, cap_text in captions:
                cap_bbox = cap_block[:4]
                cap_center_y = (cap_bbox[1] + cap_bbox[3]) / 2
                cap_center_x = (cap_bbox[0] + cap_bbox[2]) / 2

                # Prefer images in the same column as the caption (Miller's is a
                # two-column layout; matching by vertical proximity alone often
                # grabs a figure from the facing column at a similar page height).
                same_col_imgs = [img for img in images if _column_of(img['bbox']) == _column_of(cap_bbox)]
                candidate_imgs = same_col_imgs if same_col_imgs else images

                for img in candidate_imgs:
                    img_bbox = img['bbox']
                    img_center_y = (img_bbox[1] + img_bbox[3]) / 2
                    img_center_x = (img_bbox[0] + img_bbox[2]) / 2
                    # Vertical proximity dominates; horizontal offset is a tiebreaker
                    # so stacked figures within the same column don't get conflated.
                    dist = abs(img_center_y - cap_center_y) + 0.25 * abs(img_center_x - cap_center_x)
                    if dist < 400:
                        all_candidate_pairs.append((dist, cap_idx, img))

            all_candidate_pairs.sort(key=lambda t: t[0])
            assigned_image_for_caption = {}
            claimed_image_xrefs = set()
            for dist, cap_idx, img in all_candidate_pairs:
                if cap_idx in assigned_image_for_caption or img['xref'] in claimed_image_xrefs:
                    continue
                assigned_image_for_caption[cap_idx] = (img, dist)
                claimed_image_xrefs.add(img['xref'])

            for cap_idx, cap_block, cap_text in captions:
                cap_bbox = cap_block[:4]
                closest_img, min_dist = assigned_image_for_caption.get(cap_idx, (None, float('inf')))

                if closest_img and min_dist < 400:
                    xref = closest_img['xref']
                    fig_id = sanitize_fig_id(cap_text, xref, page_num)
                    
                    img_bbox = closest_img['bbox']
                    rect = fitz.Rect(img_bbox)
                    rect.x0 = max(0, rect.x0 - 5)
                    rect.y0 = max(0, rect.y0 - 5)
                    rect.x1 = min(page.rect.width, rect.x1 + 5)
                    rect.y1 = min(page.rect.height, rect.y1 + 5)
                    
                    try:
                        pix = page.get_pixmap(clip=rect, dpi=300)
                        out_filename = f"{source_file.replace('.', '_')}_{fig_id}.png"
                        out_path = os.path.join(scratch_dir, out_filename)
                        pix.save(out_path)
                        
                        heading = get_closest_heading(blocks, cap_idx)

                        # Read the PDF's own digitized text for this region instead of
                        # OCR'ing the rendered crop. This is non-scanned content, so the
                        # exact Unicode glyphs (including Greek letters/symbols that OCR
                        # frequently mangles, e.g. "β, α, and θ") are already available —
                        # OCR here would just be a lossy approximation of ground truth.
                        text_bounding_boxes = native_text_boxes_in_rect(page, rect)
                        if not text_bounding_boxes:
                            # Rare fallback: e.g. a figure with no text layer at all
                            # (pure raster photo). OCR is the only option here.
                            text_bounding_boxes = cluster_ocr_phrases(out_path)
                        # Drop page furniture (header/page number) that fell inside the crop.
                        text_bounding_boxes = strip_furniture_phrases(text_bounding_boxes, running_furniture)

                        # Relationship tracing: prefer real vector-path geometry (exact
                        # connector wires + actual arrowhead direction) over pixel
                        # sampling. Vector annotations are sometimes drawn on top of an
                        # embedded raster image (e.g. an annotated photo/scan), so this
                        # is worth trying even though the figure itself is a bitmap.
                        MAX_NETWORK_NODES = 60
                        source_target_vectors = []
                        if 1 < len(text_bounding_boxes) <= MAX_NETWORK_NODES:
                            source_target_vectors = extract_vector_relationships(
                                get_page_drawings(), rect, 300.0 / 72.0, text_bounding_boxes,
                                warnings, fig_id, page_num
                            )
                            if not source_target_vectors:
                                # Fallback: no usable vector path data here (a pure raster
                                # photo has none) — pixel-sampling is the only option left.
                                try:
                                    img = Image.open(out_path)
                                    img_rgb = img.convert('RGB')
                                    diag = math.hypot(img.width, img.height)
                                    for idx_a in range(len(text_bounding_boxes)):
                                        for idx_b in range(len(text_bounding_boxes)):
                                            if idx_a != idx_b:
                                                box_a = text_bounding_boxes[idx_a]
                                                box_b = text_bounding_boxes[idx_b]

                                                # Cheap prefilter before the pixel-sampling check:
                                                # boxes farther apart than the image diagonal can't
                                                # be connected by a line drawn inside this crop.
                                                ca = box_a['bbox']
                                                cb = box_b['bbox']
                                                cdist = math.hypot(
                                                    (ca[0] + ca[2]) / 2 - (cb[0] + cb[2]) / 2,
                                                    (ca[1] + ca[3]) / 2 - (cb[1] + cb[3]) / 2
                                                )
                                                if cdist > diag:
                                                    continue

                                                # Verify physical vector line pixels locally
                                                text_a = box_a['text'].replace(" [uncertain]", "")
                                                text_b = box_b['text'].replace(" [uncertain]", "")
                                                if not (_is_substantive_label(text_a) and _is_substantive_label(text_b)):
                                                    continue
                                                if verify_pixel_connection(img_rgb, box_a['bbox'], box_b['bbox']):
                                                    source_target_vectors.append({
                                                        "source": text_a,
                                                        "target": text_b,
                                                        "status": "locally_verified",
                                                        "direction": "Detected",
                                                        **infer_sign_and_neurotransmitter(text_a, text_b)
                                                    })
                                except Exception:
                                    pass
                        elif len(text_bounding_boxes) > MAX_NETWORK_NODES:
                            warnings.append(
                                f"Figure {fig_id} on page {page_num} has "
                                f"{len(text_bounding_boxes)} labels (> {MAX_NETWORK_NODES}); "
                                "relationship tracing was skipped for this figure."
                            )
                        
                        visual_data_engines.append({
                            "id": fig_id,
                            "sourceFile": source_file,
                            "pageNumber": page_num,
                            "image_path": os.path.abspath(out_path),
                            "closest_text_heading": heading,
                            "caption": cap_text,
                            "text_bounding_boxes": text_bounding_boxes,
                            "archetype": "ANATOMICAL NETWORKS & MICROCIRCUIT MAPS",
                            "details": {
                                "nodes": {
                                    b['text'].replace(" [uncertain]", ""): {
                                        "name": b['text'].replace(" [uncertain]", ""),
                                        "relational_variants": []
                                    } for b in text_bounding_boxes
                                },
                                "source_target_vectors": source_target_vectors
                            }
                        })
                    except Exception as e:
                        warnings.append(f"Failed to crop visual engine {fig_id} on page {page_num}: {str(e)}")
                else:
                    # Enforce programmatic Keyword-Anchor Spatial Cropping for PDF vector diagrams
                    fig_id = sanitize_fig_id(cap_text, f"vector_{cap_idx}", page_num)
                    
                    try:
                        # 1. Look for vector drawings around this caption bbox
                        # Search window above (up to 350 pts)
                        drawings = get_page_drawings()
                        cy0 = cap_bbox[1]
                        cy1 = cap_bbox[3]
                        
                        drawings_above = []
                        drawings_below = []
                        for d in drawings:
                            r = d["rect"]
                            # Exclude full-bleed page-border/rule lines only — a real
                            # waveform trace (ECG/capnography strip) is ALSO thin and
                            # often page-wide, so width/height alone can't distinguish
                            # them. A rule line is drawn as 1-2 path items; a trace is
                            # drawn as many connected segments tracing an oscillation.
                            # Require both conditions before excluding, so a thin-but-
                            # complex path (the trace) survives even if a thin-and-simple
                            # one (the rule line) still gets filtered out.
                            is_full_bleed = r.width >= page.rect.width * 0.98 or r.height >= page.rect.height * 0.98
                            is_simple_rule = len(d.get("items", [])) <= 2
                            if is_full_bleed and is_simple_rule:
                                continue
                            
                            # Clamped to page bounds
                            if cy0 - 350 <= r.y0 < cy0 and r.y1 <= cy0 + 10:
                                drawings_above.append(d)
                            elif cy1 - 10 <= r.y0 < cy1 + 350:
                                drawings_below.append(d)
                                
                        # Decide which zone has vector diagrams
                        drawings_zone = []
                        is_above = True
                        if len(drawings_above) > len(drawings_below) and len(drawings_above) > 0:
                            drawings_zone = drawings_above
                            is_above = True
                        elif len(drawings_below) >= len(drawings_above) and len(drawings_below) > 0:
                            drawings_zone = drawings_below
                            is_above = False
                            
                        if len(drawings_zone) > 0:
                            # Union box of drawings in page coordinate space
                            union_rect = fitz.Rect()
                            for d in drawings_zone:
                                r = d["rect"]
                                clamped = fitz.Rect(max(0, r.x0), max(0, r.y0), min(page.rect.width, r.x1), min(page.rect.height, r.y1))
                                union_rect.include_rect(clamped)
                                
                            # Expand horizontal bounds to align with caption or margins
                            x0 = max(45.0, min(cap_bbox[0] - 10, union_rect.x0 - 10))
                            x1 = min(page.rect.width - 45.0, max(cap_bbox[2] + 10, union_rect.x1 + 10))
                            
                            if is_above:
                                y1 = cy0 - 2
                                y0 = max(0, union_rect.y0 - 15)
                                # Avoid preceding caption overlapping
                                for other_idx, other_block, other_text in captions:
                                    if other_idx != cap_idx:
                                        other_bbox = other_block[:4]
                                        # If other caption is above this one
                                        if other_bbox[3] < cy0:
                                            y0 = max(other_bbox[3] + 2, y0)
                            else:
                                y0 = cy1 + 2
                                y1 = min(page.rect.height, union_rect.y1 + 15)
                                # Avoid succeeding caption overlapping
                                for other_idx, other_block, other_text in captions:
                                    if other_idx != cap_idx:
                                        other_bbox = other_block[:4]
                                        # If other caption is below this one
                                        if other_bbox[1] > cy1:
                                            y1 = min(other_bbox[1] - 2, y1)
                                            
                            rect = fitz.Rect(x0, y0, x1, y1)
                            
                            # Render at 300 DPI
                            pix = page.get_pixmap(clip=rect, dpi=300)
                            out_filename = f"{source_file.replace('.', '_')}_{fig_id}.png"
                            out_path = os.path.join(scratch_dir, out_filename)
                            pix.save(out_path)
                            
                            heading = get_closest_heading(blocks, cap_idx)
                            # Native PDF text for this region (axis labels, legends) is
                            # exact Unicode — prefer it over OCR for the same reason as
                            # the embedded-image branch above.
                            text_bounding_boxes = native_text_boxes_in_rect(page, rect)
                            if not text_bounding_boxes:
                                text_bounding_boxes = cluster_ocr_phrases(out_path)
                            # Drop page furniture (header/page number) that fell inside the crop.
                            text_bounding_boxes = strip_furniture_phrases(text_bounding_boxes, running_furniture)

                            # Archetype classification by caption keyword
                            archetype = "COORDINATE X-Y GRAPHS & COMPLEMENTARY PANELS"
                            details = {}
                            
                            cap_lower = cap_text.lower()
                            # Modality-neutral guess; the TS-side StrategyRouter (see
                            # archetypes.ts's CONTINUOUS_WAVEFORM config) has final say
                            # and additionally classifies which modality this is.
                            if any(kw in cap_lower for kw in (
                                "eeg", "electroencephalogram", "ecg", "ekg", "electrocardiogram",
                                "capnogr", "etco2", "plethysmogr", "waveform", "tracing"
                            )):
                                archetype = "CONTINUOUS_WAVEFORM"
                                details = {
                                    "nodes": {},
                                    "source_target_vectors": [],
                                    "labels": [b["text"].replace(" [uncertain]", "") for b in text_bounding_boxes]
                                }
                            elif "hypnogram" in cap_lower:
                                archetype = "TIMELINE_STEP_CHART_HYPNOGRAM"
                                details = {
                                    "labels": [b["text"].replace(" [uncertain]", "") for b in text_bounding_boxes],
                                    "y_axis_stages": ["W", "REM", "1", "2", "3"],
                                    "x_axis_timestamps": ["10pm", "11pm", "12pm", "01am", "02am", "03am", "04am", "05am", "06am"]
                                }
                            else:
                                # This branch only runs when get_drawings() already found
                                # vector content here (that's how this crop region was
                                # located in the first place), so it's the best case for
                                # real relationship tracing via actual connector/arrowhead
                                # geometry rather than pixel sampling.
                                vector_edges = []
                                if 1 < len(text_bounding_boxes) <= 60:
                                    vector_edges = extract_vector_relationships(
                                        drawings, rect, 300.0 / 72.0, text_bounding_boxes,
                                        warnings, fig_id, page_num
                                    )
                                if vector_edges:
                                    archetype = "ANATOMICAL NETWORKS & MICROCIRCUIT MAPS"
                                    details = {
                                        "nodes": {
                                            b['text'].replace(" [uncertain]", ""): {
                                                "name": b['text'].replace(" [uncertain]", ""),
                                                "relational_variants": []
                                            } for b in text_bounding_boxes
                                        },
                                        "source_target_vectors": vector_edges
                                    }
                                else:
                                    details = {
                                        "labels": [b["text"].replace(" [uncertain]", "") for b in text_bounding_boxes]
                                    }

                            # Attach deterministic axis/series structure for X-Y graphs
                            # (exact from the text layer; never fabricated).
                            if "COORDINATE" in (archetype or ""):
                                gs = extract_graph_structure(text_bounding_boxes)
                                if gs:
                                    details = {**details, "graph_structure": gs}

                            visual_data_engines.append({
                                "id": fig_id,
                                "sourceFile": source_file,
                                "pageNumber": page_num,
                                "image_path": os.path.abspath(out_path),
                                "closest_text_heading": heading,
                                "caption": cap_text,
                                "text_bounding_boxes": text_bounding_boxes,
                                "archetype": archetype,
                                "details": details
                            })
                    except Exception as e:
                        warnings.append(f"Failed to vector-crop visual engine {fig_id} on page {page_num}: {str(e)}")
                        
        char_count = len(raw_text)
        content_type = "empty"
        if char_count > 0:
            content_type = "text" if char_count > 50 else "mixed"
            
        fragments.append({
            "id": page_id,
            "sourceFile": source_file,
            "pageNumber": page_num,
            "contentType": content_type,
            "rawText": raw_text,
            "characterCount": char_count,
            "parsedSections": sections,
            "word_bounding_boxes": word_bounding_boxes
        })
        
    return {
        "fragments": fragments,
        "visual_data_engines": visual_data_engines,
        "warnings": warnings,
        "quality": compute_quality_metrics(fragments)
    }

def extract_image(file_path):
    source_file = os.path.basename(file_path)
    lang = 'eng_best' if os.path.exists('/opt/homebrew/share/tessdata/eng_best.traineddata') else 'eng'
    
    raw_text = ""
    warnings = []
    word_bounding_boxes = []
    
    try:
        img = Image.open(file_path)
        custom_config = f'--oem 1 --psm 6 -l {lang}'
        
        # Word coordinates OCR
        data = pytesseract.image_to_data(img, config=custom_config, output_type=pytesseract.Output.DICT)
        raw_text_parts = []
        for i in range(len(data['text'])):
            word = data['text'][i].strip()
            try:
                c = float(data['conf'][i])
            except (ValueError, TypeError):
                c = 0.0
            if word and c > 0:
                x0 = data['left'][i]
                y0 = data['top'][i]
                x1 = x0 + data['width'][i]
                y1 = y0 + data['height'][i]
                
                word_bounding_boxes.append({
                    "text": word,
                    "bbox": [x0, y0, x1, y1]
                })
                raw_text_parts.append(word)
        raw_text = " ".join(raw_text_parts)
    except Exception as e:
        warnings.append(f"Image OCR failed: {str(e)}")
        
    char_count = len(raw_text)
    
    fragments = [{
        "id": "PAGE_001",
        "sourceFile": source_file,
        "pageNumber": 1,
        "contentType": "text" if char_count > 50 else "mixed",
        "rawText": raw_text,
        "characterCount": char_count,
        "parsedSections": [{
            "heading": "",
            "body": raw_text,
            "startLine": 1,
            "category": "paragraph"
        }],
        "word_bounding_boxes": word_bounding_boxes
    }]
    
    scratch_dir = SCRATCH_DIR
    os.makedirs(scratch_dir, exist_ok=True)

    out_path = os.path.join(scratch_dir, f"{source_file.replace('.', '_')}_full.png")
    try:
        import shutil
        shutil.copy(file_path, out_path)
    except Exception:
        pass

    text_bounding_boxes = strip_furniture_phrases(cluster_ocr_phrases(out_path))
    
    visual_data_engines = [{
        "id": "FIG_IMAGE_FULL",
        "sourceFile": source_file,
        "pageNumber": 1,
        "image_path": os.path.abspath(out_path),
        "closest_text_heading": "",
        "caption": raw_text[:200] + "..." if len(raw_text) > 200 else raw_text,
        "text_bounding_boxes": text_bounding_boxes,
        "archetype": "ANATOMICAL NETWORKS & MICROCIRCUIT MAPS",
        "details": {
            "nodes": {
                b['text'].replace(" [uncertain]", ""): {
                    "name": b['text'].replace(" [uncertain]", ""),
                    "relational_variants": []
                } for b in text_bounding_boxes
            },
            "source_target_vectors": []
        }
    }]
    
    return {
        "fragments": fragments,
        "visual_data_engines": visual_data_engines,
        "warnings": warnings,
        "quality": compute_quality_metrics(fragments)
    }

def extract_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    source_file = os.path.basename(file_path)
    
    fragments = []
    for idx, slide in enumerate(prs.slides):
        slide_parts = []
        if slide.shapes.title:
            slide_parts.append(slide.shapes.title.text)
            
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text:
                        slide_parts.append(paragraph.text)
                        
        slide_text = "\n".join(slide_parts)
        char_count = len(slide_text)
        
        fragments.append({
            "id": f"PAGE_{idx+1:03d}",
            "sourceFile": source_file,
            "pageNumber": idx + 1,
            "contentType": "text" if char_count > 50 else "mixed",
            "rawText": slide_text,
            "characterCount": char_count,
            "parsedSections": [{
                "heading": slide.shapes.title.text if slide.shapes.title else "",
                "body": slide_text,
                "startLine": 1,
                "category": "paragraph"
            }]
        })
        
    return {
        "fragments": fragments,
        "visual_data_engines": [],
        "warnings": []
    }

def main():
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: python3 local_parser.py <file_path>"}))
        sys.exit(1)
        
    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(json.dumps({"error": f"File not found: {file_path}"}))
        sys.exit(1)
        
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == '.pdf':
            result = extract_pdf(file_path)
        elif ext in ('.png', '.jpg', '.jpeg'):
            result = extract_image(file_path)
        elif ext == '.pptx':
            result = extract_pptx(file_path)
        else:
            print(json.dumps({"error": f"Unsupported extension: {ext}"}))
            sys.exit(1)
            
        print(json.dumps(result))
    except Exception as e:
        print(json.dumps({
            "error": str(e),
            "traceback": traceback.format_exc()
        }))
        sys.exit(1)

if __name__ == "__main__":
    main()
